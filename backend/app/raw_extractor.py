import html
import json
import os
import platform
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any

import fitz
from fastapi import UploadFile

from app.config import get_settings


RAW_EXTENSIONS = {".docx", ".xlsx", ".pptx", ".pdf"}
OFFICE_PDF_FILTERS = {
    ".docx": "writer_pdf_Export",
    ".xlsx": "calc_pdf_Export",
    ".pptx": "impress_pdf_Export",
}
PDF_EXPORT_OPTIONS = {
    "MaxImageResolution": {"type": "long", "value": "300"},
    "Quality": {"type": "long", "value": "95"},
    "ReduceImageResolution": {"type": "boolean", "value": "false"},
    "EmbedStandardFonts": {"type": "boolean", "value": "true"},
    "EmbedFonts": {"type": "boolean", "value": "true"},
    "SubsetFonts": {"type": "boolean", "value": "false"},
}


class RawExtractionError(ValueError):
    pass


def validate_raw_upload(filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix not in RAW_EXTENSIONS:
        raise RawExtractionError("Only DOCX, XLSX, PPTX, and PDF files are supported")
    return suffix


def save_raw_upload(upload: UploadFile, raw_id: str) -> tuple[str, str, Path, int]:
    suffix = validate_raw_upload(upload.filename or "")
    settings = get_settings()
    raw_dir = settings.resolved_raw_storage_dir / raw_id
    raw_dir.mkdir(parents=True, exist_ok=True)
    original_path = raw_dir / f"original{suffix}"

    size = 0
    with original_path.open("wb") as destination:
        while True:
            chunk = upload.file.read(1024 * 1024)
            if not chunk:
                break
            size += len(chunk)
            destination.write(chunk)

    return upload.filename or original_path.name, suffix[1:], original_path, size


def create_raw_outputs(source_path: Path, source_format: str) -> tuple[Path, Path, list[str]]:
    suffix = f".{source_format.lower()}"
    output_dir = source_path.parent
    pdf_path = output_dir / "preview.pdf"
    html_path = output_dir / "content.html"
    warnings: list[str] = []

    create_pdf_preview(source_path, suffix, pdf_path)
    html_path.write_text(build_html_document(source_path, suffix), encoding="utf-8")
    if html_path.stat().st_size == 0:
        warnings.append("empty_html_output")
    return pdf_path, html_path, warnings


def create_pdf_preview(source_path: Path, suffix: str, pdf_path: Path) -> None:
    if suffix == ".pdf":
        shutil.copy2(source_path, pdf_path)
        return
    convert_office_to_pdf(source_path, suffix, pdf_path)


def convert_office_to_pdf(source_path: Path, suffix: str, pdf_path: Path) -> None:
    export_filter = OFFICE_PDF_FILTERS.get(suffix)
    if not export_filter:
        raise RawExtractionError(f"PDF conversion is not supported for {suffix}")

    soffice = find_libreoffice()
    with tempfile.TemporaryDirectory(prefix="raw2pdf_") as tmp:
        tmp_path = Path(tmp)
        cmd = [
            str(soffice),
            "--headless",
            "--convert-to",
            f"pdf:{export_filter}:{json.dumps(PDF_EXPORT_OPTIONS)}",
            "--outdir",
            str(tmp_path),
            str(source_path),
        ]
        completed = subprocess.run(cmd, capture_output=True, text=True, timeout=600, env=libreoffice_environment())
        if completed.returncode != 0:
            raise RawExtractionError(
                f"LibreOffice failed to convert the document to PDF: {completed.stderr or completed.stdout}"
            )
        produced = tmp_path / f"{source_path.stem}.pdf"
        if not produced.is_file():
            raise RawExtractionError("LibreOffice did not produce a PDF preview")
        shutil.move(str(produced), pdf_path)


def find_libreoffice() -> Path:
    settings = get_settings()
    candidates: list[Path] = []
    if settings.libreoffice_path:
        candidates.append(Path(settings.libreoffice_path))
    if platform.system() == "Darwin":
        candidates.append(Path("/Applications/LibreOffice.app/Contents/MacOS/soffice"))
    candidates.extend(
        [
            Path("/usr/bin/soffice"),
            Path("/usr/local/bin/soffice"),
            Path("/opt/homebrew/bin/soffice"),
            Path("/opt/libreoffice25.2/program/soffice"),
        ]
    )
    for candidate in candidates:
        if candidate.is_file():
            return candidate
    raise RawExtractionError(
        "LibreOffice was not found. Install LibreOffice or set LIBREOFFICE_PATH in the backend .env."
    )


def libreoffice_environment() -> dict[str, str]:
    env = os.environ.copy()
    if platform.system() == "Darwin":
        env["PYTHONPATH"] = (
            "/Applications/LibreOffice.app/Contents/Frameworks/"
            "LibreOfficePython.framework/Versions/3.10/lib/python3.10"
        )
        env["PYTHONHOME"] = "/Applications/LibreOffice.app/Contents/Frameworks/LibreOfficePython.framework/Versions/3.10"
    elif platform.system() == "Linux":
        env["PYTHONPATH"] = "/opt/libreoffice25.2/program"
        env["PYTHONHOME"] = "/opt/libreoffice25.2/program/python"
    return env


def build_html_document(source_path: Path, suffix: str) -> str:
    if suffix == ".docx":
        body = _docx_to_html(source_path)
    elif suffix == ".xlsx":
        body = _xlsx_to_html(source_path)
    elif suffix == ".pptx":
        body = _pptx_to_html(source_path)
    elif suffix == ".pdf":
        body = _pdf_to_html(source_path)
    else:
        raise RawExtractionError(f"HTML extraction is not supported for {suffix}")

    title = html.escape(source_path.name)
    return f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>{title}</title>
  <style>
    body {{ color: #1d2529; font-family: Inter, Arial, sans-serif; line-height: 1.55; margin: 24px; }}
    h1, h2, h3 {{ line-height: 1.2; }}
    section {{ border-bottom: 1px solid #dde4df; margin-bottom: 24px; padding-bottom: 18px; }}
    table {{ border-collapse: collapse; margin: 12px 0; width: 100%; }}
    td, th {{ border: 1px solid #cbd5cf; padding: 6px 8px; vertical-align: top; }}
    th {{ background: #eef3ef; }}
    pre {{ white-space: pre-wrap; }}
  </style>
</head>
<body>
  <h1>{title}</h1>
  {body}
</body>
</html>"""


def _docx_to_html(source_path: Path) -> str:
    import bleach
    import mammoth

    with source_path.open("rb") as source:
        result = mammoth.convert_to_html(source)
    allowed_tags = {
        "a",
        "b",
        "br",
        "em",
        "h1",
        "h2",
        "h3",
        "h4",
        "h5",
        "h6",
        "i",
        "li",
        "ol",
        "p",
        "strong",
        "table",
        "tbody",
        "td",
        "th",
        "thead",
        "tr",
        "u",
        "ul",
    }
    allowed_attrs = {"a": ["href"], "td": ["colspan", "rowspan"], "th": ["colspan", "rowspan"]}
    return bleach.clean(result.value, tags=allowed_tags, attributes=allowed_attrs, strip=True)


def _xlsx_to_html(source_path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(source_path, read_only=True, data_only=True)
    sections: list[str] = []
    try:
        for sheet in workbook.worksheets:
            rows = list(sheet.iter_rows(values_only=True))
            rows = [row for row in rows if any(cell is not None for cell in row)]
            table = _rows_to_table(rows)
            sections.append(f"<section><h2>{html.escape(sheet.title)}</h2>{table}</section>")
    finally:
        workbook.close()
    return "\n".join(sections) or "<p>No readable worksheet data found.</p>"


def _pptx_to_html(source_path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(source_path)
    sections: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        fragments = [f"<h2>Slide {index}</h2>"]
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = _shape_text_to_html(shape)
                if text:
                    fragments.append(text)
            if getattr(shape, "has_table", False):
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                fragments.append(_rows_to_table(rows))
        sections.append(f"<section>{''.join(fragments)}</section>")
    return "\n".join(sections) or "<p>No readable slide content found.</p>"


def _pdf_to_html(source_path: Path) -> str:
    sections: list[str] = []
    with fitz.open(source_path) as document:
        for page_index, page in enumerate(document, start=1):
            blocks = page.get_text("blocks")
            fragments = [f"<h2>Page {page_index}</h2>"]
            for block in sorted(blocks, key=lambda item: (item[1], item[0])):
                text = str(block[4]).strip()
                if text:
                    fragments.append(f"<p>{html.escape(text).replace(chr(10), '<br />')}</p>")
            sections.append(f"<section>{''.join(fragments)}</section>")
    return "\n".join(sections) or "<p>No readable PDF text found.</p>"


def _shape_text_to_html(shape: Any) -> str:
    paragraphs: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if text:
            paragraphs.append(f"<p>{html.escape(text)}</p>")
    return "".join(paragraphs)


def _rows_to_table(rows: list[Any]) -> str:
    if not rows:
        return "<p>No table data found.</p>"
    html_rows: list[str] = []
    for row_index, row in enumerate(rows):
        tag = "th" if row_index == 0 else "td"
        cells = "".join(f"<{tag}>{html.escape(_cell_to_text(cell))}</{tag}>" for cell in row)
        html_rows.append(f"<tr>{cells}</tr>")
    return f"<table><tbody>{''.join(html_rows)}</tbody></table>"


def _cell_to_text(value: Any) -> str:
    if value is None:
        return ""
    return str(value)
