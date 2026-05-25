import base64
import io
import json
import os
import threading
import time
from datetime import datetime, timedelta
from types import SimpleNamespace
import zipfile

import pytest

try:
    import pymupdf as fitz
except ImportError:  # pragma: no cover - compatibility for older PyMuPDF installs
    import fitz
from PIL import Image

from app.config import get_settings
from app.models import Document, WorkflowRun, WorkflowRunItem
from tests.conftest import get_client


ONE_BY_ONE_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
)
SCHEMA_COUNTER = 0
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


def assert_xlsx_response(response) -> None:
    assert response.status_code == 200, response.text
    assert response.headers["content-type"] == XLSX_MIME
    assert response.content.startswith(b"PK")


def test_health() -> None:
    with get_client() as client:
        response = client.get("/api/health")
    assert response.status_code == 200
    assert response.json() == {"status": "ok"}


def test_system_status_mock_mode(monkeypatch) -> None:
    try:
        monkeypatch.setenv("VLM_PROVIDER", "mock")
        monkeypatch.setenv("UPLOAD_MAX_BATCH_FILES", "1234")
        monkeypatch.setenv("UPLOAD_CHUNK_FILES", "10")
        monkeypatch.setenv("WORKFLOW_MAX_WORKERS", "16")
        monkeypatch.setenv("VLM_MAX_CONCURRENT_REQUESTS", "8")
        get_settings.cache_clear()
        with get_client() as client:
            response = client.get("/api/system/status")
        assert response.status_code == 200
        payload = response.json()
        assert payload["vlm_provider"] == "mock"
        assert payload["is_mock"] is True
        assert payload["upload_max_batch_files"] == 1234
        assert payload["upload_chunk_files"] == 10
        assert payload["workflow_max_workers"] == 16
        assert payload["vlm_max_concurrent_requests"] == 8
        assert "batch_max_workers" not in payload
        assert payload["document_page_max_long_edge"] == 3000
        assert payload["document_page_jpeg_quality"] == 88
        assert "vlm_api_key" not in payload
    finally:
        get_settings.cache_clear()


def test_database_pool_defaults_to_64_connections() -> None:
    from app.database import engine

    pool_size = getattr(engine.pool, "size", None)
    assert callable(pool_size)
    assert pool_size() == 64
    assert getattr(engine.pool, "_max_overflow", None) == 0


def test_root_env_upsert_creates_vlm_settings(monkeypatch, tmp_path) -> None:
    from app import config as config_module

    env_path = tmp_path / ".env"
    monkeypatch.setattr(config_module, "ROOT_ENV_PATH", env_path)
    config_module.upsert_root_env(
        {
            "VLM_API_KEY": "test-secret",
            "VLM_MODEL_NAME": "test-model",
            "LIBREOFFICE_PATH": "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        },
        include_defaults=True,
    )

    contents = env_path.read_text(encoding="utf-8")
    assert 'APP_ENV="local"' in contents
    assert 'VLM_API_KEY="test-secret"' in contents
    assert 'VLM_MODEL_NAME="test-model"' in contents
    assert 'LIBREOFFICE_PATH="/Applications/LibreOffice.app/Contents/MacOS/soffice"' in contents


def test_vlm_settings_include_libreoffice_path(monkeypatch, tmp_path) -> None:
    from app import config as config_module

    env_path = tmp_path / ".env"
    monkeypatch.setattr(config_module, "ROOT_ENV_PATH", env_path)

    with get_client() as client:
        response = client.put(
            "/api/settings/vlm",
            json={
                "api_key": "test-secret",
                "model_name": "test-model",
                "libreoffice_path": "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                "provider": "openai",
                "workflow_max_workers": 11,
                "vlm_max_concurrent_requests": 7,
                "kie_field_group_size": 3,
            },
        )
        assert response.status_code == 200, response.text
        payload = response.json()
        assert payload["libreoffice_path"] == "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        assert "vlm_max_concurrent_requests" in payload
        assert payload["workflow_max_workers"] == 11
        assert "batch_max_workers" not in payload
        assert "kie_field_group_size" in payload
        assert payload["runtime_settings_writable"] is True

    contents = env_path.read_text(encoding="utf-8")
    assert 'LIBREOFFICE_PATH="/Applications/LibreOffice.app/Contents/MacOS/soffice"' in contents
    assert "BATCH_MAX_WORKERS" not in contents
    assert 'WORKFLOW_MAX_WORKERS="11"' in contents
    assert 'VLM_MAX_CONCURRENT_REQUESTS="7"' in contents
    assert 'KIE_FIELD_GROUP_SIZE="3"' in contents


def test_vlm_settings_are_readonly_in_production(monkeypatch, tmp_path) -> None:
    from app import config as config_module

    env_path = tmp_path / ".env"
    monkeypatch.setattr(config_module, "ROOT_ENV_PATH", env_path)
    monkeypatch.setenv("APP_ENV", "production")
    monkeypatch.delenv("ALLOW_RUNTIME_SETTINGS", raising=False)
    get_settings.cache_clear()

    try:
        with get_client() as client:
            settings_response = client.get("/api/settings/vlm")
            assert settings_response.status_code == 200
            assert settings_response.json()["runtime_settings_writable"] is False

            response = client.put(
                "/api/settings/vlm",
                json={
                    "api_key": "test-secret",
                    "model_name": "test-model",
                    "libreoffice_path": "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                    "provider": "openai",
                },
            )
            assert response.status_code == 403
            assert not env_path.exists()
    finally:
        get_settings.cache_clear()


def test_cors_origin_env_parser() -> None:
    from app.config import parse_cors_allowed_origins

    assert parse_cors_allowed_origins("https://app.example.com, https://admin.example.com") == [
        "https://app.example.com",
        "https://admin.example.com",
    ]


def test_shared_secret_auth_and_csrf(monkeypatch) -> None:
    monkeypatch.setenv("ACCESS_CONTROL_MODE", "shared_secret")
    monkeypatch.setenv("APP_ACCESS_SECRET", "test-access-code")
    monkeypatch.setenv("SESSION_SECRET_KEY", "test-session-secret")
    monkeypatch.setenv("SESSION_COOKIE_SECURE", "false")
    get_settings.cache_clear()

    try:
        with get_client() as client:
            blocked = client.get("/api/documents")
            assert blocked.status_code == 401

            bad = client.post("/api/auth/session", json={"access_code": "wrong"})
            assert bad.status_code == 401

            created = client.post("/api/auth/session", json={"access_code": "test-access-code"})
            assert created.status_code == 200, created.text
            payload = created.json()
            assert payload["authenticated"] is True
            assert payload["csrf_token"]
            assert "httponly" in created.headers["set-cookie"].lower()

            allowed = client.get("/api/documents")
            assert allowed.status_code == 200

            missing_csrf = client.post("/api/schemas", json={"name": "blocked", "fields": []})
            assert missing_csrf.status_code == 403

            created_schema = client.post(
                "/api/schemas",
                headers={"X-CSRF-Token": payload["csrf_token"]},
                json={
                    "name": "csrf_allowed",
                    "fields": [
                        {"key_name": "field", "description": "visible field", "output_format": "string"},
                    ],
                },
            )
            assert created_schema.status_code == 200, created_schema.text
    finally:
        monkeypatch.delenv("ACCESS_CONTROL_MODE", raising=False)
        monkeypatch.delenv("APP_ACCESS_SECRET", raising=False)
        monkeypatch.delenv("SESSION_SECRET_KEY", raising=False)
        monkeypatch.delenv("SESSION_COOKIE_SECURE", raising=False)
        get_settings.cache_clear()


def test_upload_rejects_content_type_spoof() -> None:
    with get_client() as client:
        response = client.post(
            "/api/documents",
            files={"file": ("fake.png", b"not an image", "image/png")},
        )
        assert response.status_code == 415


def test_upload_rejects_file_size_over_limit(monkeypatch) -> None:
    monkeypatch.setenv("UPLOAD_MAX_FILE_BYTES", "4")
    get_settings.cache_clear()
    try:
        with get_client() as client:
            response = client.post(
                "/api/documents",
                files={"file": ("invoice.png", ONE_BY_ONE_PNG, "image/png")},
            )
            assert response.status_code == 413
    finally:
        monkeypatch.delenv("UPLOAD_MAX_FILE_BYTES", raising=False)
        get_settings.cache_clear()


def test_upload_downscales_preview_over_pixel_limit(monkeypatch) -> None:
    image = Image.new("RGB", (20, 10), (255, 255, 255))
    stream = io.BytesIO()
    image.save(stream, format="PNG")

    monkeypatch.setenv("UPLOAD_MAX_IMAGE_PIXELS", "100")
    monkeypatch.setenv("DOCUMENT_PAGE_MAX_LONG_EDGE", "0")
    get_settings.cache_clear()
    try:
        with get_client() as client:
            response = client.post(
                "/api/documents",
                files={"file": ("oversized.png", stream.getvalue(), "image/png")},
            )
            assert response.status_code == 200, response.text
            page = response.json()["pages"][0]
            assert page["width"] * page["height"] <= 100
            assert page["width"] * page["height"] <= 98
    finally:
        monkeypatch.delenv("UPLOAD_MAX_IMAGE_PIXELS", raising=False)
        monkeypatch.delenv("DOCUMENT_PAGE_MAX_LONG_EDGE", raising=False)
        get_settings.cache_clear()


def test_batch_upload_rejects_too_many_files(monkeypatch) -> None:
    monkeypatch.setenv("UPLOAD_MAX_BATCH_FILES", "1")
    get_settings.cache_clear()
    try:
        with get_client() as client:
            schema = create_schema(client, "batch_limit_schema")
            response = client.post(
                "/api/batches",
                data={"schema_id": schema["id"]},
                files=[
                    ("files", ("a.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("b.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert response.status_code == 413
    finally:
        monkeypatch.delenv("UPLOAD_MAX_BATCH_FILES", raising=False)
        get_settings.cache_clear()


def test_large_batch_multipart_allows_configured_counts_above_parser_default(monkeypatch) -> None:
    monkeypatch.setenv("UPLOAD_MAX_BATCH_FILES", "10000")
    get_settings.cache_clear()
    try:
        with get_client() as client:
            files = [
                ("files", (f"document_{index:04d}.png", ONE_BY_ONE_PNG, "image/png"))
                for index in range(1001)
            ]
            response = client.post(
                "/api/classification-batches",
                data={"classifier_id": "missing"},
                files=files,
            )
            assert response.status_code == 404, response.text
            assert response.json()["detail"] == "Document classifier not found"
    finally:
        monkeypatch.delenv("UPLOAD_MAX_BATCH_FILES", raising=False)
        get_settings.cache_clear()


def test_large_batch_multipart_returns_413_over_configured_limit(monkeypatch) -> None:
    monkeypatch.setenv("UPLOAD_MAX_BATCH_FILES", "1001")
    get_settings.cache_clear()
    try:
        with get_client() as client:
            files = [
                ("files", (f"document_{index:04d}.png", ONE_BY_ONE_PNG, "image/png"))
                for index in range(1002)
            ]
            response = client.post(
                "/api/classification-batches",
                data={"classifier_id": "missing"},
                files=files,
            )
            assert response.status_code == 413, response.text
            assert "1001" in response.json()["detail"]
    finally:
        monkeypatch.delenv("UPLOAD_MAX_BATCH_FILES", raising=False)
        get_settings.cache_clear()


def test_retention_cleanup_removes_expired_uploads(monkeypatch) -> None:
    from app.database import SessionLocal

    monkeypatch.setenv("UPLOAD_RETENTION_HOURS", "1")
    get_settings.cache_clear()
    try:
        with get_client() as client:
            document = upload_png(client)
            db = SessionLocal()
            try:
                row = db.get(Document, document["document_id"])
                assert row is not None
                storage_path = row.storage_path
                row.created_at = datetime.utcnow() - timedelta(hours=2)
                for page in row.pages:
                    page.created_at = datetime.utcnow() - timedelta(hours=2)
                db.commit()
            finally:
                db.close()

            cleaned = client.post("/api/maintenance/retention-cleanup")
            assert cleaned.status_code == 200, cleaned.text
            assert cleaned.json()["status"] == "cleaned"
            assert client.get(f"/api/documents/{document['document_id']}").status_code == 404
            assert not os.path.exists(os.path.dirname(storage_path))
    finally:
        monkeypatch.delenv("UPLOAD_RETENTION_HOURS", raising=False)
        get_settings.cache_clear()


def test_vlm_runtime_kwargs_include_speed_controls(monkeypatch) -> None:
    from app.vlm import _build_llm_kwargs

    try:
        monkeypatch.setenv("VLM_API_KEY", "test-secret")
        monkeypatch.setenv("VLM_MODEL_NAME", "test-model")
        monkeypatch.setenv("VLM_REASONING_EFFORT", "minimal")
        monkeypatch.setenv("VLM_VERBOSITY", "low")
        monkeypatch.setenv("VLM_MAX_COMPLETION_TOKENS", "1024")
        monkeypatch.setenv("VLM_TOP_P", "0.8")
        monkeypatch.setenv("VLM_SERVICE_TIER", "auto")
        get_settings.cache_clear()

        kwargs = _build_llm_kwargs()
        assert kwargs["reasoning_effort"] == "minimal"
        assert kwargs["verbosity"] == "low"
        assert kwargs["max_completion_tokens"] == 1024
        assert kwargs["top_p"] == 0.8
        assert kwargs["service_tier"] == "auto"
    finally:
        get_settings.cache_clear()


def test_kie_field_groups_run_with_bounded_parallelism(monkeypatch) -> None:
    from app import extraction as extraction_module
    from app.extraction import DocumentPageSnapshot, DocumentSnapshot
    from app.schemas import FieldDefinition

    monkeypatch.setenv("KIE_FIELD_GROUP_SIZE", "1")
    monkeypatch.setenv("WORKFLOW_MAX_WORKERS", "1")
    monkeypatch.setenv("VLM_MAX_CONCURRENT_REQUESTS", "2")
    get_settings.cache_clear()

    active = 0
    max_active = 0
    calls = 0
    lock = threading.Lock()

    def fake_extract_with_vlm(fields, image_inputs=None):
        nonlocal active, max_active, calls
        with lock:
            active += 1
            calls += 1
            max_active = max(max_active, active)
        time.sleep(0.03)
        with lock:
            active -= 1
        return {field.key_name: {"value": f"value-{field.key_name}", "page": 1, "confidence": 1, "evidence": "mock"} for field in fields}

    try:
        monkeypatch.setattr(extraction_module, "extract_with_vlm", fake_extract_with_vlm)
        fields = [
            FieldDefinition(key_name=f"field_{index}", description="test field", output_format="string")
            for index in range(4)
        ]
        document = DocumentSnapshot(
            id="doc_1",
            storage_path="/tmp/doc_1/source.png",
            pages=[DocumentPageSnapshot(page_number=1, image_path="/tmp/doc_1/page_1.png")],
        )

        values = extraction_module._extract_grouped_values(document, fields, [], "job_1")

        assert calls == 4
        assert max_active == 2
        assert set(values) == {field.key_name for field in fields}
    finally:
        get_settings.cache_clear()


def test_module_batch_workers_use_workflow_worker_limit(monkeypatch) -> None:
    from app import document_modules

    active = 0
    max_active = 0
    finalized = False
    lock = threading.Lock()

    def fake_runner(_job_id: str) -> None:
        nonlocal active, max_active
        with lock:
            active += 1
            max_active = max(max_active, active)
        time.sleep(0.03)
        with lock:
            active -= 1

    def fake_failer(_job_id: str, message: str) -> None:
        raise AssertionError(message)

    def fake_finalizer() -> None:
        nonlocal finalized
        finalized = True

    monkeypatch.setenv("WORKFLOW_MAX_WORKERS", "2")
    get_settings.cache_clear()
    try:
        document_modules._run_parallel_batch(["job_1", "job_2", "job_3", "job_4"], fake_runner, fake_failer, fake_finalizer)
        assert max_active == 2
        assert finalized is True
    finally:
        get_settings.cache_clear()


def test_extraction_batch_workers_use_workflow_worker_limit(monkeypatch) -> None:
    from app import extraction as extraction_module

    active = 0
    max_active = 0
    finalized: list[str] = []
    lock = threading.Lock()

    def fake_run_job(_job_id: str) -> None:
        nonlocal active, max_active
        with lock:
            active += 1
            max_active = max(max_active, active)
        time.sleep(0.03)
        with lock:
            active -= 1

    monkeypatch.setenv("WORKFLOW_MAX_WORKERS", "2")
    monkeypatch.setattr(extraction_module, "run_extraction_job", fake_run_job)
    monkeypatch.setattr(extraction_module, "_finalize_batch", lambda batch_id: finalized.append(batch_id))
    get_settings.cache_clear()
    try:
        extraction_module.run_batch_jobs("batch_1", ["job_1", "job_2", "job_3", "job_4"])
        assert max_active == 2
        assert finalized == ["batch_1"]
    finally:
        get_settings.cache_clear()


def test_vlm_api_style_auto_detects_google_and_base_url(monkeypatch) -> None:
    from app.vlm import resolve_vlm_api_style

    try:
        monkeypatch.setenv("VLM_PROVIDER", "auto")
        monkeypatch.setenv("VLM_API_KEY", "AIzaSyCP_test_key")
        monkeypatch.setenv("VLM_MODEL_NAME", "gemini-3.1-flash-lite")
        monkeypatch.delenv("VLM_BASE_URL", raising=False)
        get_settings.cache_clear()
        assert resolve_vlm_api_style() == "google_genai"

        monkeypatch.setenv("VLM_BASE_URL", "https://openrouter.ai/api/v1")
        get_settings.cache_clear()
        assert resolve_vlm_api_style() == "openai_compatible"

        monkeypatch.setenv("VLM_PROVIDER", "openai")
        monkeypatch.delenv("VLM_BASE_URL", raising=False)
        get_settings.cache_clear()
        assert resolve_vlm_api_style() == "google_genai"
    finally:
        get_settings.cache_clear()


def test_vlm_errors_have_stable_codes_and_redact_secrets(monkeypatch) -> None:
    from app.vlm import VlmRuntimeError, _coerce_structured_response, _sanitize_provider_error, resolve_vlm_api_style

    try:
        monkeypatch.setenv("VLM_PROVIDER", "unknown_provider")
        get_settings.cache_clear()
        with pytest.raises(VlmRuntimeError) as provider_error:
            resolve_vlm_api_style()
        assert provider_error.value.code == "VLM_PROVIDER_UNSUPPORTED"
        assert provider_error.value.as_detail()["code"] == "VLM_PROVIDER_UNSUPPORTED"

        monkeypatch.setenv("VLM_API_KEY", "AIzaSyCP_test_key_should_not_leak")
        get_settings.cache_clear()
        sanitized = _sanitize_provider_error(RuntimeError("bad key AIzaSyCP_test_key_should_not_leak"))
        assert "AIzaSyCP_test_key_should_not_leak" not in sanitized
        assert "[redacted]" in sanitized

        with pytest.raises(VlmRuntimeError) as response_error:
            _coerce_structured_response(SimpleNamespace(text="not-json"))
        assert response_error.value.code == "VLM_RESPONSE_INVALID_JSON"
    finally:
        get_settings.cache_clear()


def test_vlm_provider_request_retries_transient_broken_pipe(monkeypatch) -> None:
    from app import vlm as vlm_module
    from app.vlm import VlmRuntimeError, _invoke_vlm_with_limit

    calls = 0

    def fake_invoke(*_args, **_kwargs):
        nonlocal calls
        calls += 1
        if calls == 1:
            raise VlmRuntimeError("VLM_PROVIDER_REQUEST_FAILED", "Google GenAI VLM request failed: [Errno 32] Broken pipe")
        return {"ok": True}

    try:
        monkeypatch.setenv("VLM_MAX_RETRIES", "2")
        get_settings.cache_clear()
        monkeypatch.setattr(vlm_module.time, "sleep", lambda _delay: None)
        monkeypatch.setattr(vlm_module, "_invoke_structured_llm", fake_invoke)

        assert _invoke_vlm_with_limit("system", "prompt", [], {}, "google_genai") == {"ok": True}
        assert calls == 2
    finally:
        get_settings.cache_clear()


def test_google_generation_config_uses_structured_output_and_thinking_level(monkeypatch) -> None:
    from app.vlm import _build_google_generation_config

    try:
        monkeypatch.setenv("VLM_PROVIDER", "auto")
        monkeypatch.setenv("VLM_API_KEY", "AIzaSyCP_test_key")
        monkeypatch.setenv("VLM_MODEL_NAME", "gemini-3.1-flash-lite")
        monkeypatch.setenv("VLM_REASONING_EFFORT", "minimal")
        get_settings.cache_clear()

        schema = {"type": "object", "properties": {"name": {"type": "string"}}, "required": ["name"]}
        config = _build_google_generation_config("system", schema)
        assert config["system_instruction"] == "system"
        assert config["response_mime_type"] == "application/json"
        assert config["response_json_schema"] == schema
        assert config["thinking_config"] == {"thinking_level": "minimal"}
    finally:
        get_settings.cache_clear()


def test_classification_schema_uses_classes_or_unknown_only() -> None:
    from app.vlm import _classification_output_schema
    from app.schemas import ClassCandidate

    schema = _classification_output_schema(
        classes=[ClassCandidate(class_name="contract", description="Contract", signals=[])],
        allow_unknown=False,
    )
    assert schema["properties"]["status"]["enum"] == ["classified", "unknown"]
    class_name_schema = schema["properties"]["class_name"]
    assert {"type": "null"} in class_name_schema["anyOf"]


def test_classification_validation_clears_unknown_class_name() -> None:
    from app.document_modules import ClassificationContext, _validate_classification_output
    from app.extraction import DocumentSnapshot
    from app.schemas import ClassCandidate

    context = ClassificationContext(
        document=DocumentSnapshot(id="doc_1", storage_path="", pages=[]),
        classifier_id="clf_1",
        classes=[ClassCandidate(class_name="contract", description="Contract", signals=[])],
        allow_unknown=True,
    )
    output = _validate_classification_output(
        {
            "status": "unknown",
            "class_name": "contract",
            "confidence": 0.3,
            "reason": "No class matched.",
            "evidence": [],
        },
        context,
    )
    assert output["status"] == "unknown"
    assert output["class_name"] is None


def test_schema_validation_and_creation() -> None:
    with get_client() as client:
        invalid = client.post(
            "/api/schemas",
            json={
                "name": "bad_schema",
                "fields": [
                    {"key_name": "total", "description": "Total amount", "output_format": "float"},
                    {"key_name": "total", "description": "Duplicate total", "output_format": "float"},
                ],
            },
        )
        assert invalid.status_code == 422

        unsupported_format = client.post(
            "/api/schemas",
            json={
                "name": "bad_format",
                "fields": [
                    {"key_name": "count", "description": "Unsupported integer field", "output_format": "int"},
                ],
            },
        )
        assert unsupported_format.status_code == 422

        valid = create_schema(client)
        assert valid["name"] == "invoice_basic"
        assert valid["fields"][0]["key_name"] == "invoice_number"

        region_schema = client.post(
            "/api/schemas",
            json={
                "name": "region_schema",
                "regions": [
                    {"id": "region_1", "name": "Region 1", "page": 1, "x": 0.1, "y": 0.2, "width": 0.3, "height": 0.1}
                ],
                "fields": [
                    {
                        "key_name": "handwritten_name",
                        "description": "손글씨 이름 영역",
                        "output_format": "string",
                        "region_id": "region_1",
                    },
                ],
            },
        )
        assert region_schema.status_code == 200, region_schema.text
        assert region_schema.json()["regions"][0]["x"] == 0.1
        assert region_schema.json()["fields"][0]["region_id"] == "region_1"

        invalid_region = client.post(
            "/api/schemas",
            json={
                "name": "invalid_region_schema",
                "regions": [
                    {"id": "region_1", "name": "Region 1", "page": 1, "x": 0.8, "y": 0.2, "width": 0.3, "height": 0.1}
                ],
                "fields": [
                    {
                        "key_name": "handwritten_name",
                        "description": "손글씨 이름 영역",
                        "output_format": "string",
                        "region_id": "region_1",
                    },
                ],
            },
        )
        assert invalid_region.status_code == 422

        korean_with_space = client.post(
            "/api/schemas",
            json={
                "name": "korean_schema",
                "fields": [
                    {"key_name": "법정 대리인 성", "description": "우측 하단의 법정 대리인 성명", "output_format": "string"},
                ],
            },
        )
        assert korean_with_space.status_code == 200
        assert korean_with_space.json()["fields"][0]["key_name"] == "법정 대리인 성"

        screenshot_payload = client.post(
            "/api/schemas",
            json={
                "name": "document_schema",
                "fields": [
                    {"key_name": "개정일", "description": "좌측 하단의 개정일자", "output_format": "date"},
                    {"key_name": "본인 성명", "description": "우측 하단의 본인 성명", "output_format": "string"},
                ],
            },
        )
        assert screenshot_payload.status_code == 200, screenshot_payload.text
        assert screenshot_payload.json()["fields"][1]["key_name"] == "본인 성명"


def test_image_upload() -> None:
    with get_client() as client:
        document = upload_png(client)
        assert document["page_count"] == 1
        assert document["created_at"]
        image = client.get(document["pages"][0]["image_url"])
        assert image.status_code == 200
        assert image.headers["content-type"] == "image/jpeg"
        thumbnail = client.get(f"/api/documents/{document['document_id']}/pages/1/thumbnail?width=96")
        assert thumbnail.status_code == 200
        assert thumbnail.headers["content-type"] == "image/jpeg"

        documents = client.get("/api/documents").json()
        assert any(item["document_id"] == document["document_id"] for item in documents)


def test_jpeg_upload_preserves_source_pixels_with_dpi_metadata() -> None:
    image = Image.new("RGB", (300, 420), (255, 255, 255))
    buffer = io.BytesIO()
    image.save(buffer, format="JPEG", dpi=(300, 300))

    with get_client() as client:
        response = client.post(
            "/api/documents",
            files={"file": ("scan.jpg", buffer.getvalue(), "image/jpeg")},
        )
        assert response.status_code == 200, response.text
        document = response.json()
        assert document["pages"][0]["width"] == 300
        assert document["pages"][0]["height"] == 420
        image_response = client.get(document["pages"][0]["image_url"])
        assert image_response.status_code == 200
        loaded = Image.open(io.BytesIO(image_response.content))
        assert loaded.size == (300, 420)


def test_pdf_upload() -> None:
    with get_client() as client:
        pdf_bytes = make_pdf_bytes()
        response = client.post(
            "/api/documents",
            files={"file": ("invoice.pdf", pdf_bytes, "application/pdf")},
        )
        assert response.status_code == 200
        document = response.json()
        assert document["page_count"] == 1
        assert document["pages"][0]["width"] > 0
        assert document["pages"][0]["height"] > 0


def test_office_upload_for_key_information_extractor(monkeypatch) -> None:
    def fake_convert(source_path, suffix, pdf_path):
        document = fitz.open()
        page = document.new_page(width=240, height=120)
        page.insert_text((24, 60), f"Converted {source_path.name}")
        document.save(pdf_path)
        document.close()

    monkeypatch.setattr("app.document_processor.convert_office_to_pdf", fake_convert)

    samples = [
        ("report.docx", make_docx_bytes(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
        ("deck.pptx", make_pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
    ]
    with get_client() as client:
        for filename, data, mime_type in samples:
            response = client.post(
                "/api/documents",
                files={"file": (filename, data, mime_type)},
            )
            assert response.status_code == 200, response.text
            document = response.json()
            assert document["filename"] == filename
            assert document["page_count"] == 1
            image = client.get(document["pages"][0]["image_url"])
            assert image.status_code == 200
            assert image.headers["content-type"] == "image/jpeg"


def test_extraction_fails_without_vlm_credentials(monkeypatch) -> None:
    monkeypatch.setenv("VLM_API_KEY", "")
    monkeypatch.setenv("VLM_MODEL_NAME", "")
    monkeypatch.setenv("OPENAI_API_KEY", "")
    monkeypatch.setenv("OPENAI_MODEL_NAME", "")
    get_settings.cache_clear()
    with get_client() as client:
        document = upload_png(client)
        schema = create_schema(client)
        job_response = client.post(
            "/api/extraction-jobs",
            json={"document_id": document["document_id"], "schema_id": schema["id"]},
        )
        assert job_response.status_code == 200
        job_id = job_response.json()["job_id"]

        job = client.get(f"/api/extraction-jobs/{job_id}").json()
        assert job["status"] == "failed"
        assert job["result_id"] is None
        assert "VLM API key and model name are required" in job["error_message"]


def test_schema_update_replaces_current_schema() -> None:
    with get_client() as client:
        schema = create_schema(client)
        updated = client.patch(
            f"/api/schemas/{schema['id']}",
            json={
                "display_name": "Updated Invoice Basic",
                "fields": [
                    {
                        "key_name": "invoice_number",
                        "description": "Invoice number near the top.",
                        "output_format": "string",
                    },
                    {
                        "key_name": "invoice_date",
                        "description": "Invoice issue date.",
                        "output_format": "date",
                    },
                ],
            },
        )
        assert updated.status_code == 200, updated.text
        payload = updated.json()
        assert "current_version" not in payload
        assert payload["display_name"] == "Updated Invoice Basic"
        assert payload["fields"][1]["key_name"] == "invoice_date"


def test_schema_update_allows_same_name_for_loaded_schema() -> None:
    with get_client() as client:
        schema = create_schema(client, name="테스트")
        updated = client.patch(
            f"/api/schemas/{schema['id']}",
            json={
                "name": "테스트",
                "display_name": "테스트",
                "description": "수정된 설명",
                "fields": [
                    {
                        "key_name": "수정필드",
                        "description": "사용자가 저장된 스키마를 불러온 뒤 수정한 필드",
                        "output_format": "string",
                    }
                ],
            },
        )
        assert updated.status_code == 200, updated.text
        payload = updated.json()
        assert payload["id"] == schema["id"]
        assert payload["name"] == "테스트"
        assert payload["description"] == "수정된 설명"
        assert "current_version" not in payload
        assert payload["fields"][0]["key_name"] == "수정필드"


def test_schema_update_merges_duplicate_loaded_schema_name() -> None:
    from app.database import SessionLocal
    from app.models import Schema

    with get_client() as client:
        schema = create_schema(client, name="중복스키마")
        duplicate_schema_json = {
            "name": "중복스키마",
            "display_name": "중복스키마",
            "description": "old duplicate",
            "is_template": False,
            "template_category": None,
            "pinned": False,
            "regions": [],
            "fields": [
                {
                    "key_name": "old_field",
                    "description": "Old duplicate field.",
                    "output_format": "string",
                }
            ],
        }
        db = SessionLocal()
        try:
            duplicate = Schema(
                name="중복스키마",
                display_name="중복스키마",
                description="old duplicate",
                current_version=1,
                schema_json=json.dumps(duplicate_schema_json, ensure_ascii=False),
                is_template=False,
                template_category=None,
                pinned=False,
                ephemeral=False,
            )
            db.add(duplicate)
            db.commit()
            duplicate_id = duplicate.id
        finally:
            db.close()

        updated = client.patch(
            f"/api/schemas/{schema['id']}",
            json={
                "name": "중복스키마",
                "display_name": "중복스키마",
                "description": "merged current schema",
                "fields": [
                    {
                        "key_name": "current_field",
                        "description": "Current schema field.",
                        "output_format": "string",
                    }
                ],
            },
        )
        assert updated.status_code == 200, updated.text
        schemas = [item for item in client.get("/api/schemas").json() if item["name"] == "중복스키마"]
        assert len(schemas) == 1
        assert schemas[0]["id"] == schema["id"]
        assert client.get(f"/api/schemas/{duplicate_id}").status_code == 404


def test_schema_delete_archives_and_allows_name_reuse() -> None:
    with get_client() as client:
        schema = create_schema(client, name="삭제테스트")
        deleted = client.delete(f"/api/schemas/{schema['id']}")
        assert deleted.status_code == 200, deleted.text
        assert deleted.json()["archived"] is True

        listed_names = [item["name"] for item in client.get("/api/schemas").json()]
        assert "삭제테스트" not in listed_names

        archived = client.get(f"/api/schemas/{schema['id']}")
        assert archived.status_code == 200
        assert archived.json()["archived"] is True

        recreated = create_schema(client, name="삭제테스트")
        assert recreated["id"] != schema["id"]


def test_schema_duplicate_copies_definition_with_incremented_name() -> None:
    with get_client() as client:
        schema = create_schema(client, name="성명")
        duplicated = client.post(f"/api/schemas/{schema['id']}/duplicate")
        assert duplicated.status_code == 200, duplicated.text
        payload = duplicated.json()
        assert payload["id"] != schema["id"]
        assert payload["name"] == "성명 (1)"
        assert payload["display_name"] == "성명 (1)"
        assert payload["fields"] == schema["fields"]

        duplicated_again = client.post(f"/api/schemas/{schema['id']}/duplicate")
        assert duplicated_again.status_code == 200, duplicated_again.text
        assert duplicated_again.json()["name"] == "성명 (2)"


def test_schema_recommendation_mock_mode() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            document = upload_png(client)
            response = client.post("/api/schemas/recommendations", json={"document_id": document["document_id"]})
            assert response.status_code == 200, response.text
            payload = response.json()
            assert payload["name"] == "ai_recommended_schema"
            assert len(payload["fields"]) >= 3
            assert {field["output_format"] for field in payload["fields"]} <= {"string", "float", "date", "bool"}
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_required_field_checklist_recommendation_mock_mode() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            document = upload_png(client)
            response = client.post(
                "/api/required-field-checklists/recommendations",
                json={"document_id": document["document_id"]},
            )
            assert response.status_code == 200, response.text
            payload = response.json()
            assert payload["name"] == "ai_recommended_checklist"
            assert len(payload["items"]) >= 3
            assert {item["evidence_type"] for item in payload["items"]} <= {
                "text_or_handwriting",
                "checkbox",
                "signature_or_stamp",
                "visual_mark",
                "other",
            }
            assert {item["region_id"] for item in payload["items"] if item["region_id"]} <= {
                region["id"] for region in payload["regions"]
            }
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_schema_description_recommendation_mock_mode() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            response = client.post(
                "/api/schemas/description-recommendations",
                json={
                    "name": "consent_schema",
                    "current_description": "Old description",
                    "fields": [
                        {
                            "key_name": "본인 성명",
                            "description": "문서 하단 서명 영역의 본인 성명",
                            "output_format": "string",
                        },
                        {
                            "key_name": "동의 여부",
                            "description": "체크박스 선택 상태를 기준으로 한 동의 여부",
                            "output_format": "bool",
                        },
                    ],
                },
            )
            assert response.status_code == 200, response.text
            payload = response.json()
            assert "consent_schema" in payload["description"]
            assert "본인 성명" in payload["description"]
            assert payload["reasoning"]
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_extraction_mock_mode_returns_evidence_and_normalized_values() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            document = upload_png(client)
            schema = create_schema(client)
            job_response = client.post(
                "/api/extraction-jobs",
                json={"document_id": document["document_id"], "schema_id": schema["id"]},
            )
            assert job_response.status_code == 200, job_response.text
            job_id = job_response.json()["job_id"]

            job = client.get(f"/api/extraction-jobs/{job_id}").json()
            assert job["status"] == "completed"
            values = job["result"]["validated_output"]["values"]
            assert values["invoice_number"]["page"] == 1
            assert values["invoice_number"]["evidence"]
            assert values["total_amount"]["normalized_value"] == 1234.5

            jobs = client.get(f"/api/extraction-jobs?document_id={document['document_id']}").json()
            assert any(item["job_id"] == job_id for item in jobs)

            csv_export = client.get(f"/api/extraction-results/{job['result_id']}/export?format=csv")
            assert csv_export.status_code == 200
            assert csv_export.content.startswith(b"\xef\xbb\xbf")
            assert "charset=utf-8" in csv_export.headers["content-type"]
            assert "evidence" in csv_export.text.splitlines()[0]

            xlsx_export = client.get(f"/api/extraction-results/{job['result_id']}/export?format=xlsx")
            assert_xlsx_response(xlsx_export)

            corrected_output = job["result"]["validated_output"]
            corrected_output["values"]["invoice_number"]["value"] = "INV-EDITED"
            patch = client.patch(
                f"/api/extraction-results/{job['result_id']}",
                json={"corrected_output": corrected_output},
            )
            assert patch.status_code == 200, patch.text
            assert patch.json()["corrected_output"]["values"]["invoice_number"]["value"] == "INV-EDITED"
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_draft_extraction_does_not_list_schema() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            document = upload_png(client)
            response = client.post(
                "/api/extraction-jobs/draft",
                json={
                    "document_id": document["document_id"],
                    "schema": {
                        "name": "unsaved_draft_schema",
                        "display_name": "Unsaved Draft Schema",
                        "fields": [
                            {
                                "key_name": "draft_value",
                                "description": "Value visible in the draft document.",
                                "output_format": "string",
                            }
                        ],
                    },
                },
            )
            assert response.status_code == 200, response.text
            job = client.get(f"/api/extraction-jobs/{response.json()['job_id']}").json()
            assert job["status"] == "completed"
            assert job["result"]["validated_output"]["values"]["draft_value"]["value"] == "Sample draft_value"

            schemas = client.get("/api/schemas").json()
            assert all(item["name"] != "unsaved_draft_schema" for item in schemas)
            hidden_schema = client.get(f"/api/schemas/{job['schema_id']}")
            assert hidden_schema.status_code == 200
            assert hidden_schema.json()["ephemeral"] is True
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_schema_name_conflict_and_clear_parsing_history() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="conflict_schema")
            duplicate = client.post(
                "/api/schemas",
                json={
                    "name": "conflict_schema",
                    "fields": [
                        {
                            "key_name": "other",
                            "description": "Other field.",
                            "output_format": "string",
                        }
                    ],
                },
            )
            assert duplicate.status_code == 409

            document = upload_png(client)
            job_response = client.post(
                "/api/extraction-jobs",
                json={"document_id": document["document_id"], "schema_id": schema["id"]},
            )
            assert job_response.status_code == 200
            job = client.get(f"/api/extraction-jobs/{job_response.json()['job_id']}").json()
            assert job["status"] == "completed"

            cleared = client.delete("/api/maintenance/parsing-history")
            assert cleared.status_code == 200, cleared.text
            payload = cleared.json()
            assert payload["status"] == "cleared"
            assert payload["counts"]["documents"] >= 1
            assert client.get("/api/documents").json() == []
            assert client.get("/api/extraction-jobs").json() == []
            assert any(item["id"] == schema["id"] for item in client.get("/api/schemas").json())
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_extraction_releases_db_connection_during_vlm_call(monkeypatch) -> None:
    from app.database import engine

    checked_out_counts: list[int] = []

    def fake_extract(fields, image_paths=None, image_inputs=None):
        checkedout = getattr(engine.pool, "checkedout", None)
        checked_out_counts.append(checkedout() if checkedout else 0)
        return {
            "invoice_number": {"value": "INV-POOL-001", "page": 1, "evidence": "test", "confidence": 0.9},
            "total_amount": {"value": "10.00", "page": 1, "evidence": "test", "confidence": 0.9},
        }

    monkeypatch.setattr("app.extraction.extract_with_vlm", fake_extract)

    with get_client() as client:
        document = upload_png(client)
        schema = create_schema(client)
        response = client.post(
            "/api/extraction-jobs",
            json={"document_id": document["document_id"], "schema_id": schema["id"]},
        )
        assert response.status_code == 200, response.text
        job_id = response.json()["job_id"]
        job = client.get(f"/api/extraction-jobs/{job_id}").json()
        assert job["status"] == "completed"

    assert checked_out_counts
    assert max(checked_out_counts) == 0


def test_extraction_uses_schema_regions_for_cropped_inputs(monkeypatch) -> None:
    captured_calls: list[dict[str, object]] = []

    def fake_extract(fields, image_paths=None, image_inputs=None):
        captured_calls.append({"fields": fields, "image_paths": image_paths, "image_inputs": image_inputs})
        field_names = {field.key_name for field in fields}
        values = {}
        if "handwritten_name" in field_names:
            values["handwritten_name"] = {"value": "홍길동", "page": 1, "evidence": "region crop", "confidence": 0.91}
        if "handwritten_phone" in field_names:
            values["handwritten_phone"] = {"value": "010-0000-0000", "page": 1, "evidence": "region crop", "confidence": 0.9}
        if "document_date" in field_names:
            values["document_date"] = {"value": "2026.05.20", "page": 1, "evidence": "full page", "confidence": 0.9}
        return {
            **values,
        }

    monkeypatch.setattr("app.extraction.extract_with_vlm", fake_extract)

    with get_client() as client:
        document = client.post(
            "/api/documents",
            files={"file": ("sample.pdf", make_pdf_bytes(), "application/pdf")},
        ).json()
        schema = client.post(
            "/api/schemas",
            json={
                "name": "mixed_region_schema",
                "regions": [
                    {"id": "region_1", "name": "Handwriting block", "page": 1, "x": 0.1, "y": 0.1, "width": 0.5, "height": 0.4}
                ],
                "fields": [
                    {
                        "key_name": "handwritten_name",
                        "description": "손글씨 이름 영역",
                        "output_format": "string",
                        "region_id": "region_1",
                    },
                    {
                        "key_name": "handwritten_phone",
                        "description": "손글씨 연락처 영역",
                        "output_format": "string",
                        "region_id": "region_1",
                    },
                    {
                        "key_name": "document_date",
                        "description": "문서 전체에서 날짜",
                        "output_format": "date",
                    },
                ],
            },
        ).json()
        job_response = client.post(
            "/api/extraction-jobs",
            json={"document_id": document["document_id"], "schema_id": schema["id"]},
        )
        assert job_response.status_code == 200, job_response.text
        job = client.get(f"/api/extraction-jobs/{job_response.json()['job_id']}").json()

    assert job["status"] == "completed"
    assert len(captured_calls) == 2

    full_page_call = next(call for call in captured_calls if [field.key_name for field in call["fields"]] == ["document_date"])
    full_page_inputs = full_page_call["image_inputs"]
    assert isinstance(full_page_inputs, list)
    assert any("Full document page 1" in item["label"] for item in full_page_inputs)

    region_call = next(call for call in captured_calls if {field.key_name for field in call["fields"]} == {"handwritten_name", "handwritten_phone"})
    region_inputs = region_call["image_inputs"]
    assert isinstance(region_inputs, list)
    assert any("Full page context" in item["label"] and "Handwriting block" in item["label"] for item in region_inputs)
    assert any("Masked full page context" in item["label"] and "Handwriting block" in item["label"] for item in region_inputs)
    assert any("Cropped extraction region" in item["label"] and "Handwriting block" in item["label"] for item in region_inputs)
    assert any("handwritten_name, handwritten_phone" in item["label"] for item in region_inputs)
    assert len(region_inputs) == 3
    region_fields = region_call["fields"]
    assert region_fields[0].region_id == "region_1"
    assert region_fields[1].region_id == "region_1"


def test_kie_prompts_are_split_by_region_presence() -> None:
    from app.prompts.kie import build_extraction_prompt, build_region_judgement_prompt
    from app.schemas import FieldDefinition

    full_prompt = build_extraction_prompt([
        FieldDefinition(key_name="document_date", description="Date visible on the document.", output_format="date")
    ])
    region_field = FieldDefinition(key_name="signature", description="Signature in the lower right area.", output_format="string", region_id="region_1")
    region_prompt = build_extraction_prompt([region_field])
    judgement_prompt = build_region_judgement_prompt(region_field, "서명", "first-stage evidence")

    assert "full document page images" in full_prompt
    assert "user-designated extraction region" not in full_prompt
    assert "labeled extraction region images" in region_prompt
    assert "crop image is already the user-designated extraction region" in region_prompt
    assert "verification step, not a re-extraction step" in judgement_prompt
    assert "default decision is judgement_status=correct" in judgement_prompt
    assert "If text such as 성명, 서명, 법정" in judgement_prompt


def test_kie_ai_judgement_corrects_enabled_field(monkeypatch) -> None:
    def fake_extract(fields, image_paths=None, image_inputs=None):
        return {
            field.key_name: {
                "value": "WRONG" if field.key_name == "invoice_number" else "1,234.50",
                "page": 1,
                "evidence": f"first evidence {field.key_name}",
                "confidence": 0.82,
            }
            for field in fields
        }

    judgement_calls = []
    correction_calls = []

    def fake_judge(field, initial_value, initial_evidence, image_inputs):
        judgement_calls.append((field.key_name, initial_value, initial_evidence, image_inputs))
        return {
            "judgement_status": "needs_correction",
            "reason": "The visible invoice number is different from the first-stage value.",
            "confidence": 0.93,
            "evidence": "Invoice number area",
        }

    def fake_correct(field, initial_value, initial_evidence, judgement_reason, image_inputs):
        correction_calls.append((field.key_name, judgement_reason, image_inputs))
        return {
            "value": "INV-2026-001",
            "page": 1,
            "evidence": "Corrected invoice number",
            "confidence": 0.95,
            "correction_reason": "Correct value is visible near the top.",
        }

    monkeypatch.setattr("app.extraction.extract_with_vlm", fake_extract)
    monkeypatch.setattr("app.extraction.judge_extraction_with_vlm", fake_judge)
    monkeypatch.setattr("app.extraction.correct_extraction_with_vlm", fake_correct)

    with get_client() as client:
        document = upload_png(client)
        schema = client.post(
            "/api/schemas",
            json={
                "name": "judgement_schema",
                "fields": [
                    {
                        "key_name": "invoice_number",
                        "description": "Invoice number near the top.",
                        "output_format": "string",
                        "judgement_enabled": True,
                    },
                    {
                        "key_name": "total_amount",
                        "description": "Final total amount.",
                        "output_format": "float",
                    },
                ],
            },
        )
        assert schema.status_code == 200, schema.text
        duplicated = client.post(f"/api/schemas/{schema.json()['id']}/duplicate")
        assert duplicated.status_code == 200, duplicated.text
        assert duplicated.json()["fields"][0]["judgement_enabled"] is True

        job_response = client.post(
            "/api/extraction-jobs",
            json={"document_id": document["document_id"], "schema_id": schema.json()["id"]},
        )
        assert job_response.status_code == 200, job_response.text
        job = client.get(f"/api/extraction-jobs/{job_response.json()['job_id']}").json()

    assert job["status"] == "completed"
    values = job["result"]["validated_output"]["values"]
    assert values["invoice_number"]["value"] == "INV-2026-001"
    assert values["invoice_number"]["ai_review"]["corrected"] is True
    assert values["invoice_number"]["ai_review"]["initial_value"] == "WRONG"
    assert "ai_review" not in values["total_amount"]
    assert job["result"]["raw_model_output"]["invoice_number"]["value"] == "WRONG"
    assert [call[0] for call in judgement_calls] == ["invoice_number"]
    assert [call[0] for call in correction_calls] == ["invoice_number"]


def test_kie_ai_judgement_failure_needs_review_without_failing_job(monkeypatch) -> None:
    def fake_extract(fields, image_paths=None, image_inputs=None):
        return {
            field.key_name: {
                "value": "INV-1",
                "page": 1,
                "evidence": "first evidence",
                "confidence": 0.82,
            }
            for field in fields
        }

    def fake_judge(field, initial_value, initial_evidence, image_inputs):
        raise RuntimeError("synthetic judgement failure")

    monkeypatch.setattr("app.extraction.extract_with_vlm", fake_extract)
    monkeypatch.setattr("app.extraction.judge_extraction_with_vlm", fake_judge)

    with get_client() as client:
        document = upload_png(client)
        schema = client.post(
            "/api/schemas",
            json={
                "name": "judgement_failure_schema",
                "fields": [
                    {
                        "key_name": "invoice_number",
                        "description": "Invoice number near the top.",
                        "output_format": "string",
                        "judgement_enabled": True,
                    }
                ],
            },
        )
        assert schema.status_code == 200, schema.text
        job_response = client.post(
            "/api/extraction-jobs",
            json={"document_id": document["document_id"], "schema_id": schema.json()["id"]},
        )
        assert job_response.status_code == 200, job_response.text
        job = client.get(f"/api/extraction-jobs/{job_response.json()['job_id']}").json()

    assert job["status"] == "needs_review"
    value = job["result"]["validated_output"]["values"]["invoice_number"]
    assert value["value"] == "INV-1"
    assert "ai_review_failed" in value["warnings"]
    assert value["ai_review"]["judgement_status"] == "failed"


@pytest.mark.parametrize(
    ("initial_value", "corrected_value", "expected_warning"),
    [
        ("성명", None, "ai_correction_discarded_null"),
        ("문어", "문이", "ai_correction_large_change"),
    ],
)
def test_kie_ai_correction_preserves_risky_korean_values(
    monkeypatch,
    initial_value: str,
    corrected_value: str | None,
    expected_warning: str,
) -> None:
    def fake_extract(fields, image_paths=None, image_inputs=None):
        return {
            field.key_name: {
                "value": initial_value,
                "page": 1,
                "evidence": f"first evidence {initial_value}",
                "confidence": 0.91,
            }
            for field in fields
        }

    def fake_judge(field, initial_value, initial_evidence, image_inputs):
        return {
            "judgement_status": "needs_correction",
            "reason": "Synthetic judgement requested a risky correction.",
            "confidence": 0.98,
            "evidence": "Synthetic evidence",
        }

    def fake_correct(field, initial_value, initial_evidence, judgement_reason, image_inputs):
        return {
            "value": corrected_value,
            "page": 1,
            "evidence": "Synthetic correction evidence",
            "confidence": 0.99,
            "correction_reason": "Synthetic risky correction.",
        }

    monkeypatch.setattr("app.extraction.extract_with_vlm", fake_extract)
    monkeypatch.setattr("app.extraction.judge_extraction_with_vlm", fake_judge)
    monkeypatch.setattr("app.extraction.correct_extraction_with_vlm", fake_correct)

    with get_client() as client:
        document = upload_png(client)
        schema = client.post(
            "/api/schemas",
            json={
                "name": f"judgement_risky_{expected_warning}",
                "fields": [
                    {
                        "key_name": "법정대리인성명",
                        "description": "Legal representative name in the handwritten field.",
                        "output_format": "string",
                        "judgement_enabled": True,
                    }
                ],
            },
        )
        assert schema.status_code == 200, schema.text
        job_response = client.post(
            "/api/extraction-jobs",
            json={"document_id": document["document_id"], "schema_id": schema.json()["id"]},
        )
        assert job_response.status_code == 200, job_response.text
        job = client.get(f"/api/extraction-jobs/{job_response.json()['job_id']}").json()

    assert job["status"] == "needs_review"
    value = job["result"]["validated_output"]["values"]["법정대리인성명"]
    assert value["value"] == initial_value
    assert expected_warning in value["warnings"]
    assert value["ai_review"]["judgement_status"] == "needs_correction"
    assert value["ai_review"]["corrected"] is False


def test_batch_cancel_marks_queued_jobs_canceled(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_batch_jobs", lambda batch_id, job_ids: None)

    with get_client() as client:
        schema = create_schema(client)
        response = client.post(
            "/api/batches",
            data={"schema_id": schema["id"]},
            files=[
                ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert response.status_code == 200, response.text
        batch = response.json()
        assert batch["status"] == "running"
        assert batch["progress"] == 0

        canceled = client.post(f"/api/batches/{batch['id']}/cancel")
        assert canceled.status_code == 200, canceled.text
        payload = canceled.json()
        assert payload["status"] == "canceled"
        assert payload["canceled_count"] == 2
        assert payload["progress"] == 1
        assert payload["completed_at"] is not None
        assert {item["status"] for item in payload["items"]} == {"canceled"}


def test_batch_init_items_start_flow(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_batch_jobs", lambda batch_id, job_ids: None)

    with get_client() as client:
        schema = create_schema(client)
        initialized = client.post("/api/batches/init", json={"schema_id": schema["id"], "total_count": 2})
        assert initialized.status_code == 200, initialized.text
        batch = initialized.json()
        assert batch["status"] == "uploading"
        assert batch["uploaded_count"] == 0

        early_start = client.post(f"/api/batches/{batch['id']}/start")
        assert early_start.status_code == 422
        assert early_start.json()["detail"]["uploaded_count"] == 0

        partial = client.post(
            f"/api/batches/{batch['id']}/items",
            data={"client_file_ids": ["0:first.png:1:1"]},
            files=[("files", ("first.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert partial.status_code == 200, partial.text
        partial_resume = client.post(f"/api/batches/{batch['id']}/resume")
        assert partial_resume.status_code == 422, partial_resume.text
        assert partial_resume.json()["detail"]["uploaded_count"] == 1
        assert partial_resume.json()["detail"]["total_count"] == 2

        appended = client.post(
            f"/api/batches/{batch['id']}/items",
            data={"client_file_ids": ["1:second.png:1:1"]},
            files=[("files", ("second.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert appended.status_code == 200, appended.text
        payload = appended.json()
        assert payload["uploaded_count"] == 2
        assert payload["queued_count"] == 2
        assert payload["preprocessing_count"] == 0

        summary = client.get(f"/api/batches/{batch['id']}/summary")
        assert summary.status_code == 200
        assert summary.json()["items"] == []
        assert summary.json()["queued_count"] == 2

        started = client.post(f"/api/batches/{batch['id']}/start")
        assert started.status_code == 200, started.text
        assert started.json()["status"] == "running"


def test_batch_export_csv_and_json_mock_mode() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client)
            response = client.post(
                "/api/batches",
                data={"schema_id": schema["id"]},
                files=[
                    ("files", ("z_last.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("a_first.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert response.status_code == 200, response.text
            batch = response.json()
            assert [item["filename"] for item in batch["items"]] == ["a_first.png", "z_last.png"]

            csv_response = client.get(f"/api/batches/{batch['id']}/export?format=csv")
            assert csv_response.status_code == 200, csv_response.text
            assert csv_response.content.startswith(b"\xef\xbb\xbf")
            assert "charset=utf-8" in csv_response.headers["content-type"]
            csv_text = csv_response.text
            header = csv_text.splitlines()[0]
            assert "filename,document_id,job_id,status,error_message,invoice_number" in header
            assert "invoice_number_original" in header
            assert "invoice_number_ai_review_status" in header
            assert "total_amount_ai_corrected" in header
            assert "a_first.png" in csv_text
            assert "Sample invoice_number" in csv_text
            assert csv_text.index("a_first.png") < csv_text.index("z_last.png")

            json_response = client.get(f"/api/batches/{batch['id']}/export?format=json")
            assert json_response.status_code == 200, json_response.text
            payload = json_response.json()
            assert payload["batch_id"] == batch["id"]
            assert len(payload["rows"]) == 2
            assert [row["filename"] for row in payload["rows"]] == ["a_first.png", "z_last.png"]
            assert payload["rows"][0]["invoice_number"] == "Sample invoice_number"

            xlsx_response = client.get(f"/api/batches/{batch['id']}/export?format=xlsx")
            assert_xlsx_response(xlsx_response)
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_batch_finalizes_after_mock_jobs() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client)
            response = client.post(
                "/api/batches",
                data={"schema_id": schema["id"]},
                files=[
                    ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert response.status_code == 200, response.text
            batch = response.json()

            loaded = client.get(f"/api/batches/{batch['id']}")
            assert loaded.status_code == 200, loaded.text
            payload = loaded.json()
            assert payload["status"] == "completed"
            assert payload["progress"] == 1
            assert payload["completed_count"] == 2
            assert payload["completed_at"] is not None
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_batch_high_concurrent_request_count_does_not_exhaust_db_pool() -> None:
    previous_workers = os.environ.get("VLM_MAX_CONCURRENT_REQUESTS")
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        os.environ["VLM_MAX_CONCURRENT_REQUESTS"] = "16"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client)
            response = client.post(
                "/api/batches",
                data={"schema_id": schema["id"]},
                files=[
                    ("files", (f"batch_{index}.png", ONE_BY_ONE_PNG, "image/png"))
                    for index in range(20)
                ],
            )
            assert response.status_code == 200, response.text
            batch = response.json()

            loaded = client.get(f"/api/batches/{batch['id']}")
            assert loaded.status_code == 200, loaded.text
            payload = loaded.json()
            assert payload["status"] == "completed"
            assert payload["completed_count"] == 20
            assert payload["failed_count"] == 0

            recent = client.get("/api/batches?limit=12")
            assert recent.status_code == 200, recent.text
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        if previous_workers is None:
            os.environ.pop("VLM_MAX_CONCURRENT_REQUESTS", None)
        else:
            os.environ["VLM_MAX_CONCURRENT_REQUESTS"] = previous_workers
        get_settings.cache_clear()


def test_document_classifier_config_single_job_and_patch() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            classifier = create_document_classifier(client)
            updated = client.patch(
                f"/api/document-classifiers/{classifier['id']}",
                json={
                    "description": "Updated classifier description",
                    "allow_unknown": True,
                },
            )
            assert updated.status_code == 200, updated.text
            assert updated.json()["description"] == "Updated classifier description"

            document = upload_png(client)
            response = client.post(
                "/api/classification-jobs",
                json={"document_id": document["document_id"], "classifier_id": classifier["id"]},
            )
            assert response.status_code == 200, response.text
            job = client.get(f"/api/classification-jobs/{response.json()['job_id']}").json()
            assert job["status"] == "completed"
            output = job["result"]["validated_output"]
            assert output["status"] == "classified"
            assert output["class_name"] == "contract"
            assert output["confidence"] == 0.88

            corrected = {**output, "status": "unknown", "class_name": None}
            patch = client.patch(
                f"/api/classification-results/{job['result_id']}",
                json={"corrected_output": corrected, "reviewed": True},
            )
            assert patch.status_code == 200, patch.text
            assert patch.json()["corrected_output"]["status"] == "unknown"

            deleted = client.delete(f"/api/document-classifiers/{classifier['id']}")
            assert deleted.status_code == 200
            assert deleted.json()["archived"] is True
            assert all(item["id"] != classifier["id"] for item in client.get("/api/document-classifiers").json())
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_document_classifier_duplicate_copies_classes() -> None:
    with get_client() as client:
        classifier = create_document_classifier(client, name="신청서 분류")
        duplicated = client.post(f"/api/document-classifiers/{classifier['id']}/duplicate")
        assert duplicated.status_code == 200, duplicated.text
        payload = duplicated.json()
        assert payload["id"] != classifier["id"]
        assert payload["name"] == "신청서 분류 (1)"
        assert payload["classes"] == classifier["classes"]
        assert payload["allow_unknown"] == classifier["allow_unknown"]


def test_document_classifier_batch_cancel_and_export(monkeypatch) -> None:
    with monkeypatch.context() as patch_context:
        patch_context.setattr("app.main.run_classification_batch", lambda batch_id, job_ids: None)

        with get_client() as client:
            classifier = create_document_classifier(client)
            response = client.post(
                "/api/classification-batches",
                data={"classifier_id": classifier["id"]},
                files=[
                    ("files", ("z_last.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("a_first.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert response.status_code == 200, response.text
            batch = response.json()
            assert [item["filename"] for item in batch["items"]] == ["a_first.png", "z_last.png"]

            canceled = client.post(f"/api/classification-batches/{batch['id']}/cancel")
            assert canceled.status_code == 200, canceled.text
            assert canceled.json()["status"] == "canceled"
            assert canceled.json()["canceled_count"] == 2
            assert canceled.json()["completed_at"] is not None

    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            classifier = create_document_classifier(client, name="batch_classifier")
            response = client.post(
                "/api/classification-batches",
                data={"classifier_id": classifier["id"]},
                files=[
                    ("files", ("z_last.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("a_first.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert response.status_code == 200, response.text
            batch = response.json()
            loaded = client.get(f"/api/classification-batches/{batch['id']}").json()
            assert loaded["status"] == "completed"
            assert loaded["completed_count"] == 2

            csv_response = client.get(f"/api/classification-batches/{batch['id']}/export?format=csv")
            assert csv_response.status_code == 200, csv_response.text
            assert csv_response.content.startswith(b"\xef\xbb\xbf")
            assert "charset=utf-8" in csv_response.headers["content-type"]
            assert "classification_status,class_name,confidence,reason,evidence" in csv_response.text.splitlines()[0]
            assert csv_response.text.index("a_first.png") < csv_response.text.index("z_last.png")

            json_response = client.get(f"/api/classification-batches/{batch['id']}/export?format=json")
            assert json_response.status_code == 200
            rows = json_response.json()["rows"]
            assert [row["filename"] for row in rows] == ["a_first.png", "z_last.png"]
            assert rows[0]["class_name"] == "contract"

            xlsx_response = client.get(f"/api/classification-batches/{batch['id']}/export?format=xlsx")
            assert_xlsx_response(xlsx_response)
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_required_field_checklist_single_job_and_region_validation() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            invalid = client.post(
                "/api/required-field-checklists",
                json={
                    "name": "invalid_checklist",
                    "regions": [],
                    "items": [
                        {
                            "item_name": "서명",
                            "description": "서명 존재 여부",
                            "evidence_type": "signature_or_stamp",
                            "required": True,
                            "region_id": "missing_region",
                        }
                    ],
                },
            )
            assert invalid.status_code == 422

            custom = client.post(
                "/api/required-field-checklists",
                json={
                    "name": "custom_evidence_checklist",
                    "regions": [],
                    "items": [
                        {
                            "item_name": "수기 메모",
                            "description": "사용자가 직접 정의한 증거 유형을 확인합니다.",
                            "evidence_type": "수기 메모/특이사항",
                            "required": True,
                        }
                    ],
                },
            )
            assert custom.status_code == 200, custom.text
            assert custom.json()["items"][0]["evidence_type"] == "수기 메모/특이사항"

            checklist = create_required_field_checklist(client)
            document = upload_png(client)
            response = client.post(
                "/api/required-field-check-jobs",
                json={"document_id": document["document_id"], "checklist_id": checklist["id"]},
            )
            assert response.status_code == 200, response.text
            job = client.get(f"/api/required-field-check-jobs/{response.json()['job_id']}").json()
            assert job["status"] == "needs_review"
            output = job["result"]["validated_output"]
            assert output["overall_status"] == "needs_review"
            assert [item["item_name"] for item in output["items"]] == ["성명", "서명", "체크박스"]
            assert output["items"][0]["status"] == "present"
            assert output["items"][2]["status"] == "uncertain"

            corrected = {
                **output,
                "overall_status": "complete",
                "items": [{**item, "status": "present"} for item in output["items"]],
            }
            patch = client.patch(
                f"/api/required-field-check-results/{job['result_id']}",
                json={"corrected_output": corrected, "reviewed": True},
            )
            assert patch.status_code == 200, patch.text
            assert patch.json()["corrected_output"]["overall_status"] == "complete"

            deleted = client.delete(f"/api/required-field-checklists/{checklist['id']}")
            assert deleted.status_code == 200
            assert deleted.json()["archived"] is True
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_required_checker_splits_full_page_and_region_requests(monkeypatch) -> None:
    captured_calls: list[dict[str, object]] = []

    def fake_check(items, regions, image_paths=None, image_inputs=None):
        captured_calls.append(
            {
                "item_names": [item.item_name for item in items],
                "region_ids": [region.id for region in regions],
                "labels": [item["label"] for item in image_inputs or []],
            }
        )
        return {
            "overall_status": "complete",
            "items": [
                {
                    "item_name": item.item_name,
                    "status": "present",
                    "confidence": 0.9,
                    "evidence": "captured",
                    "page": 1,
                }
                for item in items
            ],
        }

    monkeypatch.setattr("app.document_modules.check_required_fields_with_vlm", fake_check)

    with get_client() as client:
        checklist = create_required_field_checklist(client, name="split_required_checklist")
        document = upload_png(client)
        response = client.post(
            "/api/required-field-check-jobs",
            json={"document_id": document["document_id"], "checklist_id": checklist["id"]},
        )
        assert response.status_code == 200, response.text
        job = client.get(f"/api/required-field-check-jobs/{response.json()['job_id']}").json()

    assert job["status"] == "completed"
    assert len(captured_calls) == 2
    full_page_call = next(call for call in captured_calls if call["region_ids"] == [])
    assert full_page_call["item_names"] == ["성명", "체크박스"]
    assert any("Full document page" in label for label in full_page_call["labels"])

    region_call = next(call for call in captured_calls if call["region_ids"] == ["signature_region"])
    assert region_call["item_names"] == ["서명"]
    assert any("Masked context" in label and "서명 영역" in label for label in region_call["labels"])
    assert any("Cropped required field region" in label and "서명 영역" in label for label in region_call["labels"])
    assert not any("Full document page" in label for label in region_call["labels"])


def test_required_field_checklist_duplicate_copies_items_and_regions() -> None:
    with get_client() as client:
        checklist = create_required_field_checklist(client, name="필수 정보")
        duplicated = client.post(f"/api/required-field-checklists/{checklist['id']}/duplicate")
        assert duplicated.status_code == 200, duplicated.text
        payload = duplicated.json()
        assert payload["id"] != checklist["id"]
        assert payload["name"] == "필수 정보 (1)"
        assert payload["items"] == checklist["items"]
        assert payload["regions"] == checklist["regions"]


def test_required_field_check_batch_export_mock_mode() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            checklist = create_required_field_checklist(client, name="batch_checklist")
            response = client.post(
                "/api/required-field-check-batches",
                data={"checklist_id": checklist["id"]},
                files=[
                    ("files", ("z_last.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("a_first.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert response.status_code == 200, response.text
            batch = response.json()
            assert [item["filename"] for item in batch["items"]] == ["a_first.png", "z_last.png"]

            loaded = client.get(f"/api/required-field-check-batches/{batch['id']}").json()
            assert loaded["status"] == "completed"
            assert loaded["completed_count"] == 2

            csv_response = client.get(f"/api/required-field-check-batches/{batch['id']}/export?format=csv")
            assert csv_response.status_code == 200, csv_response.text
            assert csv_response.content.startswith(b"\xef\xbb\xbf")
            assert "charset=utf-8" in csv_response.headers["content-type"]
            header = csv_response.text.splitlines()[0]
            assert "overall_status" in header
            assert "성명_status" in header
            assert csv_response.text.index("a_first.png") < csv_response.text.index("z_last.png")

            json_response = client.get(f"/api/required-field-check-batches/{batch['id']}/export?format=json")
            assert json_response.status_code == 200
            rows = json_response.json()["rows"]
            assert rows[0]["성명_status"] == "present"

            xlsx_response = client.get(f"/api/required-field-check-batches/{batch['id']}/export?format=xlsx")
            assert_xlsx_response(xlsx_response)
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_required_field_check_batch_cancel_marks_queued_jobs_canceled(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_required_field_check_batch", lambda batch_id, job_ids: None)

    with get_client() as client:
        checklist = create_required_field_checklist(client, name="cancel_checklist")
        response = client.post(
            "/api/required-field-check-batches",
            data={"checklist_id": checklist["id"]},
            files=[
                ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert response.status_code == 200, response.text
        batch = response.json()
        assert batch["status"] == "running"

        canceled = client.post(f"/api/required-field-check-batches/{batch['id']}/cancel")
        assert canceled.status_code == 200, canceled.text
        payload = canceled.json()
        assert payload["status"] == "canceled"
        assert payload["canceled_count"] == 2
        assert payload["progress"] == 1
        assert payload["completed_at"] is not None
        assert {item["status"] for item in payload["items"]} == {"canceled"}


def test_workflow_definition_validation_and_branch_run_mock_mode() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="workflow_invoice_schema")
            classifier = create_document_classifier(client, name="workflow_classifier")
            checklist = create_required_field_checklist(client, name="workflow_checklist")

            invalid = client.post(
                "/api/workflows",
                json={
                    "name": "invalid_workflow",
                    "definition": {
                        "nodes": [
                            {"id": "input", "data": {"kind": "input"}},
                            {"id": "classifier", "data": {"kind": "classifier", "config": {}}},
                            {"id": "export", "data": {"kind": "export"}},
                        ],
                        "edges": [
                            {"id": "e1", "source": "input", "target": "classifier"},
                            {"id": "e2", "source": "classifier", "target": "export"},
                        ],
                    },
                },
            )
            assert invalid.status_code == 422
            assert "Classifier node classifier" in invalid.text

            workflow = client.post(
                "/api/workflows",
                json={
                    "name": "분기 워크플로우",
                    "description": "Classifier 결과에 따라 계약서 경로를 실행합니다.",
                    "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"]),
                },
            )
            assert workflow.status_code == 200, workflow.text
            assert workflow.json()["validation_warnings"] == []

            run_response = client.post(
                f"/api/workflows/{workflow.json()['id']}/runs",
                files=[
                    ("files", ("z_last.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("a_first.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert run_response.status_code == 200, run_response.text
            run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
            assert run["status"] in {"completed", "needs_review"}
            assert run["total_count"] == 2
            assert [item["filename"] for item in run["items"]] == ["a_first.png", "z_last.png"]
            first = run["items"][0]["result"]
            assert first["classification"]["class_name"] == "contract"
            assert first["branch_path"] == "class:contract"
            assert "invoice_number" in first["kie_values"]
            assert first["required_items"]["성명"]["status"] == "present"

            csv_response = client.get(f"/api/workflow-runs/{run['id']}/export?format=csv")
            assert csv_response.status_code == 200, csv_response.text
            assert csv_response.content.startswith(b"\xef\xbb\xbf")
            assert "charset=utf-8" in csv_response.headers["content-type"]
            header = csv_response.text.splitlines()[0]
            assert "classification_status" in header
            assert "upload_duration_ms" in header
            assert "inference_duration_ms" in header
            assert "kie_invoice_number" in header
            assert "required_성명_status" in header
            assert csv_response.text.index("a_first.png") < csv_response.text.index("z_last.png")

            json_response = client.get(f"/api/workflow-runs/{run['id']}/export?format=json")
            assert json_response.status_code == 200
            assert isinstance(run["upload_duration_ms"], int)
            assert isinstance(run["inference_duration_ms"], int)
            assert all(isinstance(item["upload_duration_ms"], int) for item in run["items"])
            assert all(isinstance(item["inference_duration_ms"], int) for item in run["items"])
            assert json_response.json()["rows"][0]["upload_duration_ms"] is not None
            assert json_response.json()["rows"][0]["class_name"] == "contract"

            xlsx_response = client.get(f"/api/workflow-runs/{run['id']}/export?format=xlsx")
            assert_xlsx_response(xlsx_response)
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_workflow_resume_and_discard_partial_upload(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_resume_schema")
        classifier = create_document_classifier(client, name="workflow_resume_classifier")
        checklist = create_required_field_checklist(client, name="workflow_resume_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "resume_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        initialized = client.post(f"/api/workflows/{workflow.json()['id']}/runs/init", json={"total_count": 2})
        assert initialized.status_code == 200, initialized.text
        run_id = initialized.json()["id"]

        appended = client.post(
            f"/api/workflow-runs/{run_id}/items",
            data={"client_file_ids": ["0:first.png:1:1"]},
            files=[("files", ("first.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert appended.status_code == 200, appended.text
        document_id = appended.json()["items"][0]["document_id"]
        assert appended.json()["uploaded_count"] == 1
        assert appended.json()["total_count"] == 2

        resumed = client.post(f"/api/workflow-runs/{run_id}/resume")
        assert resumed.status_code == 422, resumed.text
        payload = resumed.json()["detail"]
        assert payload["uploaded_count"] == 1
        assert payload["total_count"] == 2
        assert payload["missing_count"] == 1

        paused = client.post(f"/api/workflow-runs/{run_id}/pause")
        assert paused.status_code == 200, paused.text
        assert paused.json()["status"] == "paused"

        appended = client.post(
            f"/api/workflow-runs/{run_id}/items",
            data={"client_file_ids": ["1:second.png:1:1"]},
            files=[("files", ("second.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert appended.status_code == 200, appended.text
        started = client.post(f"/api/workflow-runs/{run_id}/start")
        assert started.status_code == 200, started.text
        payload = started.json()
        assert payload["status"] == "running"
        assert payload["total_count"] == 2
        assert payload["queued_count"] == 2

        discarded = client.post(f"/api/workflow-runs/{run_id}/discard")
        assert discarded.status_code == 200, discarded.text
        payload = discarded.json()
        assert payload["status"] == "canceled"
        assert payload["items"] == []
        assert payload["uploaded_count"] == 0
        assert client.get(f"/api/documents/{document_id}").status_code == 404


def test_workflow_items_preserve_upload_index_order(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_upload_index_schema")
        classifier = create_document_classifier(client, name="workflow_upload_index_classifier")
        checklist = create_required_field_checklist(client, name="workflow_upload_index_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "upload_index_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        initialized = client.post(f"/api/workflows/{workflow.json()['id']}/runs/init", json={"total_count": 2})
        assert initialized.status_code == 200, initialized.text
        run_id = initialized.json()["id"]

        appended = client.post(
            f"/api/workflow-runs/{run_id}/items",
            data={
                "client_file_ids": ["1:z_last.png:1:1", "0:a_first.png:1:1"],
                "upload_indexes": ["1", "0"],
            },
            files=[
                ("files", ("z_last.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("a_first.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert appended.status_code == 200, appended.text
        payload = appended.json()
        assert [item["filename"] for item in payload["items"]] == ["a_first.png", "z_last.png"]
        assert [item["upload_index"] for item in payload["items"]] == [0, 1]


def test_workflow_resume_skips_existing_upload_indexes_when_client_ids_change(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_resume_index_schema")
        classifier = create_document_classifier(client, name="workflow_resume_index_classifier")
        checklist = create_required_field_checklist(client, name="workflow_resume_index_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "resume_index_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        initialized = client.post(f"/api/workflows/{workflow.json()['id']}/runs/init", json={"total_count": 3})
        assert initialized.status_code == 200, initialized.text
        run_id = initialized.json()["id"]

        first = client.post(
            f"/api/workflow-runs/{run_id}/items",
            data={"client_file_ids": ["original:0", "original:1"], "upload_indexes": ["0", "1"]},
            files=[
                ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert first.status_code == 200, first.text
        assert first.json()["uploaded_count"] == 2

        resumed = client.post(
            f"/api/workflow-runs/{run_id}/items",
            data={
                "client_file_ids": ["changed:0", "changed:1", "changed:2"],
                "upload_indexes": ["0", "1", "2"],
            },
            files=[
                ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("third.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert resumed.status_code == 200, resumed.text
        payload = resumed.json()
        assert payload["uploaded_count"] == 3
        assert [item["upload_index"] for item in payload["items"]] == [0, 1, 2]


def test_workflow_restart_seals_missing_upload_slots(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_seal_schema")
        classifier = create_document_classifier(client, name="workflow_seal_classifier")
        checklist = create_required_field_checklist(client, name="workflow_seal_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "seal_missing_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        initialized = client.post(f"/api/workflows/{workflow.json()['id']}/runs/init", json={"total_count": 3})
        assert initialized.status_code == 200, initialized.text
        run_id = initialized.json()["id"]

        appended = client.post(
            f"/api/workflow-runs/{run_id}/items",
            data={"client_file_ids": ["original:0"], "upload_indexes": ["0"]},
            files=[("files", ("first.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert appended.status_code == 200, appended.text
        paused = client.post(f"/api/workflow-runs/{run_id}/pause")
        assert paused.status_code == 200, paused.text

        restarted = client.post(f"/api/workflow-runs/{run_id}/restart")
        assert restarted.status_code == 200, restarted.text
        payload = restarted.json()
        assert payload["id"] != run_id
        assert payload["restarted_from_run_id"] == run_id
        assert payload["status"] == "running"
        assert payload["uploaded_count"] == 3
        assert payload["failed_count"] == 2
        assert payload["queued_count"] == 1
        assert sorted(item["upload_index"] for item in payload["items"]) == [0, 1, 2]
        assert sum(1 for item in payload["items"] if item["filename"].startswith("missing_upload_")) == 2


def test_workflow_pause_preserves_uploads_and_restart_requeues(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_pause_schema")
        classifier = create_document_classifier(client, name="workflow_pause_classifier")
        checklist = create_required_field_checklist(client, name="workflow_pause_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "pause_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text
        initialized = client.post(f"/api/workflows/{workflow.json()['id']}/runs/init", json={"total_count": 2})
        assert initialized.status_code == 200, initialized.text
        run_id = initialized.json()["id"]
        appended = client.post(
            f"/api/workflow-runs/{run_id}/items",
            data={"client_file_ids": ["0:first.png:1:1", "1:second.png:1:1"], "upload_indexes": ["0", "1"]},
            files=[
                ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert appended.status_code == 200, appended.text
        started = client.post(f"/api/workflow-runs/{run_id}/start")
        assert started.status_code == 200, started.text

        paused = client.post(f"/api/workflow-runs/{run_id}/pause")
        assert paused.status_code == 200, paused.text
        payload = paused.json()
        assert payload["status"] == "paused"
        assert payload["uploaded_count"] == 2
        assert {item["status"] for item in payload["items"]} == {"paused"}

        restarted = client.post(f"/api/workflow-runs/{run_id}/restart")
        assert restarted.status_code == 200, restarted.text
        payload = restarted.json()
        assert payload["id"] != run_id
        assert payload["restarted_from_run_id"] == run_id
        assert payload["status"] == "running"
        assert payload["queued_count"] == 2
        assert {item["status"] for item in payload["items"]} == {"queued"}


def test_workflow_restart_can_create_new_run_with_current_workflow_without_copying_documents(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_restart_target_schema")
        classifier = create_document_classifier(client, name="workflow_restart_target_classifier")
        checklist = create_required_field_checklist(client, name="workflow_restart_target_checklist")
        workflow_one = client.post(
            "/api/workflows",
            json={"name": "workflow_one", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        workflow_two = client.post(
            "/api/workflows",
            json={"name": "workflow_two", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow_one.status_code == 200, workflow_one.text
        assert workflow_two.status_code == 200, workflow_two.text

        run_response = client.post(
            f"/api/workflows/{workflow_one.json()['id']}/runs",
            files=[
                ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert run_response.status_code == 200, run_response.text
        source_run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
        source_document_ids = [item["document_id"] for item in source_run["items"]]
        assert source_run["started_at"] is not None
        assert source_run["completed_at"] is None

        restarted = client.post(
            f"/api/workflow-runs/{source_run['id']}/restart",
            json={"workflow_id": workflow_two.json()["id"]},
        )
        assert restarted.status_code == 200, restarted.text
        new_run = restarted.json()
        assert new_run["id"] != source_run["id"]
        assert new_run["workflow_id"] == workflow_two.json()["id"]
        assert new_run["workflow_name"] == "workflow_two"
        assert new_run["restarted_from_run_id"] == source_run["id"]
        assert [item["document_id"] for item in new_run["items"]] == source_document_ids
        assert new_run["started_at"] is not None
        assert new_run["completed_at"] is None

        replaced_source = client.get(f"/api/workflow-runs/{source_run['id']}").json()
        assert replaced_source["status"] == "canceled"
        assert replaced_source["completed_at"] is not None


def test_workflow_enqueue_creates_waiting_run_without_copying_documents(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_enqueue_schema")
        classifier = create_document_classifier(client, name="workflow_enqueue_classifier")
        checklist = create_required_field_checklist(client, name="workflow_enqueue_checklist")
        workflow_one = client.post(
            "/api/workflows",
            json={"name": "queue_workflow_one", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        workflow_two = client.post(
            "/api/workflows",
            json={"name": "queue_workflow_two", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow_one.status_code == 200, workflow_one.text
        assert workflow_two.status_code == 200, workflow_two.text

        run_response = client.post(
            f"/api/workflows/{workflow_one.json()['id']}/runs",
            files=[
                ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert run_response.status_code == 200, run_response.text
        source_run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
        source_document_ids = [item["document_id"] for item in source_run["items"]]
        from app.database import SessionLocal

        db = SessionLocal()
        try:
            document_count_before = db.query(Document).count()
        finally:
            db.close()

        queued = client.post(
            f"/api/workflow-runs/{source_run['id']}/enqueue",
            json={"workflow_id": workflow_two.json()["id"]},
        )
        assert queued.status_code == 200, queued.text
        payload = queued.json()
        assert payload["id"] != source_run["id"]
        assert payload["workflow_id"] == workflow_two.json()["id"]
        assert payload["workflow_name"] == "queue_workflow_two"
        assert payload["status"] == "waiting"
        assert payload["workflow_run_group_id"] == source_run["id"]
        assert payload["queued_from_run_id"] == source_run["id"]
        assert payload["queue_order"] == 2
        assert [item["document_id"] for item in payload["items"]] == source_document_ids
        assert {item["status"] for item in payload["items"]} == {"queued"}
        db = SessionLocal()
        try:
            assert db.query(Document).count() == document_count_before
        finally:
            db.close()

        refreshed_source = client.get(f"/api/workflow-runs/{source_run['id']}").json()
        assert refreshed_source["workflow_run_group_id"] == source_run["id"]
        assert refreshed_source["queue_order"] == 1


def test_workflow_waiting_run_cannot_enqueue_or_start_out_of_order(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_queue_order_schema")
        classifier = create_document_classifier(client, name="workflow_queue_order_classifier")
        checklist = create_required_field_checklist(client, name="workflow_queue_order_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "queue_order_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        run_response = client.post(
            f"/api/workflows/{workflow.json()['id']}/runs",
            files=[("files", ("first.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert run_response.status_code == 200, run_response.text
        source_run_id = run_response.json()["id"]

        first_waiting = client.post(f"/api/workflow-runs/{source_run_id}/enqueue")
        second_waiting = client.post(f"/api/workflow-runs/{source_run_id}/enqueue")
        assert first_waiting.status_code == 200, first_waiting.text
        assert second_waiting.status_code == 200, second_waiting.text
        first_waiting_id = first_waiting.json()["id"]
        second_waiting_id = second_waiting.json()["id"]

        reenqueued = client.post(f"/api/workflow-runs/{first_waiting_id}/enqueue")
        assert reenqueued.status_code == 409, reenqueued.text

        blocked_by_source = client.post(f"/api/workflow-runs/{first_waiting_id}/start")
        assert blocked_by_source.status_code == 409, blocked_by_source.text

        from app.database import SessionLocal

        db = SessionLocal()
        try:
            source_run = db.get(WorkflowRun, source_run_id)
            assert source_run is not None
            source_run.status = "completed"
            source_run.completed_at = datetime.utcnow()
            for item in source_run.items:
                item.status = "completed"
                item.error_message = None
                item.completed_at = datetime.utcnow()
                item.result_json = json.dumps({"document_id": item.document_id, "filename": item.filename, "node_results": {}})
            db.commit()
        finally:
            db.close()

        out_of_order = client.post(f"/api/workflow-runs/{second_waiting_id}/start")
        assert out_of_order.status_code == 409, out_of_order.text

        started = client.post(f"/api/workflow-runs/{first_waiting_id}/start")
        assert started.status_code == 200, started.text
        assert started.json()["status"] == "running"


def test_workflow_finalize_starts_next_waiting_run(monkeypatch) -> None:
    from app import workflows as workflow_module
    from app.database import SessionLocal

    dispatched: list[tuple[str, int]] = []
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)
    monkeypatch.setattr("app.workflows._dispatch_workflow_run_async", lambda run_id, generation: dispatched.append((run_id, generation)))

    with get_client() as client:
        schema = create_schema(client, name="workflow_queue_advance_schema")
        classifier = create_document_classifier(client, name="workflow_queue_advance_classifier")
        checklist = create_required_field_checklist(client, name="workflow_queue_advance_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "queue_advance_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        run_response = client.post(
            f"/api/workflows/{workflow.json()['id']}/runs",
            files=[
                ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
            ],
        )
        assert run_response.status_code == 200, run_response.text
        source_run_id = run_response.json()["id"]
        queued = client.post(f"/api/workflow-runs/{source_run_id}/enqueue")
        assert queued.status_code == 200, queued.text
        queued_run_id = queued.json()["id"]

        db = SessionLocal()
        try:
            source_run = db.get(WorkflowRun, source_run_id)
            assert source_run is not None
            generation = source_run.execution_generation
            source_run.status = "running"
            for item in source_run.items:
                item.status = "completed"
                item.error_message = None
                item.completed_at = datetime.utcnow()
                item.result_json = json.dumps({"document_id": item.document_id, "filename": item.filename, "node_results": {}})
            db.commit()
        finally:
            db.close()

        workflow_module._finalize_workflow_run(source_run_id, generation)

        payload = client.get(f"/api/workflow-runs/{queued_run_id}").json()
        assert payload["status"] == "running"
        assert payload["queued_count"] == 2
        assert dispatched == [(queued_run_id, 1)]

        db = SessionLocal()
        try:
            queued_items = db.query(WorkflowRunItem).filter(WorkflowRunItem.run_id == queued_run_id).all()
            assert {item.execution_generation for item in queued_items} == {1}
        finally:
            db.close()


@pytest.mark.parametrize("blocked_status", ["paused", "canceled"])
def test_workflow_finalize_does_not_advance_queue_for_blocking_statuses(monkeypatch, blocked_status: str) -> None:
    from app import workflows as workflow_module
    from app.database import SessionLocal

    dispatched: list[tuple[str, int]] = []
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)
    monkeypatch.setattr("app.workflows._dispatch_workflow_run_async", lambda run_id, generation: dispatched.append((run_id, generation)))

    with get_client() as client:
        schema = create_schema(client, name=f"workflow_queue_blocked_{blocked_status}_schema")
        classifier = create_document_classifier(client, name=f"workflow_queue_blocked_{blocked_status}_classifier")
        checklist = create_required_field_checklist(client, name=f"workflow_queue_blocked_{blocked_status}_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": f"queue_blocked_{blocked_status}_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        run_response = client.post(
            f"/api/workflows/{workflow.json()['id']}/runs",
            files=[("files", ("first.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert run_response.status_code == 200, run_response.text
        source_run_id = run_response.json()["id"]
        queued = client.post(f"/api/workflow-runs/{source_run_id}/enqueue")
        assert queued.status_code == 200, queued.text
        queued_run_id = queued.json()["id"]

        db = SessionLocal()
        try:
            source_run = db.get(WorkflowRun, source_run_id)
            assert source_run is not None
            generation = source_run.execution_generation
            source_run.status = blocked_status
            for item in source_run.items:
                item.status = blocked_status
                item.error_message = f"{blocked_status} by test"
                item.completed_at = datetime.utcnow()
            db.commit()
        finally:
            db.close()

        workflow_module._finalize_workflow_run(source_run_id, generation)

        payload = client.get(f"/api/workflow-runs/{queued_run_id}").json()
        assert payload["status"] == "waiting"
        assert dispatched == []


def test_workflow_cancel_waiting_run_preserves_shared_documents(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_cancel_waiting_schema")
        classifier = create_document_classifier(client, name="workflow_cancel_waiting_classifier")
        checklist = create_required_field_checklist(client, name="workflow_cancel_waiting_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "cancel_waiting_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        run_response = client.post(
            f"/api/workflows/{workflow.json()['id']}/runs",
            files=[("files", ("first.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert run_response.status_code == 200, run_response.text
        source_run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
        source_document_id = source_run["items"][0]["document_id"]
        from app.database import SessionLocal

        db = SessionLocal()
        try:
            source_row = db.get(WorkflowRun, source_run["id"])
            assert source_row is not None
            source_row.status = "completed"
            source_row.completed_at = datetime.utcnow()
            source_row.inference_started_at = None
            for item in source_row.items:
                item.status = "completed"
                item.completed_at = source_row.completed_at
            db.commit()
        finally:
            db.close()

        queued = client.post(f"/api/workflow-runs/{source_run['id']}/enqueue")
        assert queued.status_code == 200, queued.text
        queued_run_id = queued.json()["id"]

        canceled = client.post(f"/api/workflow-runs/{queued_run_id}/cancel-waiting")
        assert canceled.status_code == 200, canceled.text
        payload = canceled.json()
        assert payload["status"] == "canceled"
        assert payload["items"][0]["status"] == "canceled"
        assert client.get(f"/api/documents/{source_document_id}").status_code == 200


def test_workflow_queue_entry_delete_removes_run_without_deleting_documents(monkeypatch) -> None:
    monkeypatch.setattr("app.main.run_workflow_run", lambda *args: None)

    with get_client() as client:
        schema = create_schema(client, name="workflow_delete_queue_schema")
        classifier = create_document_classifier(client, name="workflow_delete_queue_classifier")
        checklist = create_required_field_checklist(client, name="workflow_delete_queue_checklist")
        workflow = client.post(
            "/api/workflows",
            json={"name": "delete_queue_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
        )
        assert workflow.status_code == 200, workflow.text

        run_response = client.post(
            f"/api/workflows/{workflow.json()['id']}/runs",
            files=[("files", ("first.png", ONE_BY_ONE_PNG, "image/png"))],
        )
        assert run_response.status_code == 200, run_response.text
        source_run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
        source_document_id = source_run["items"][0]["document_id"]
        from app.database import SessionLocal

        db = SessionLocal()
        try:
            source_row = db.get(WorkflowRun, source_run["id"])
            assert source_row is not None
            source_row.status = "completed"
            source_row.completed_at = datetime.utcnow()
            source_row.inference_started_at = None
            for item in source_row.items:
                item.status = "completed"
                item.completed_at = source_row.completed_at
            db.commit()
        finally:
            db.close()

        queued = client.post(f"/api/workflow-runs/{source_run['id']}/enqueue")
        assert queued.status_code == 200, queued.text
        queued_run_id = queued.json()["id"]

        removed = client.delete(f"/api/workflow-runs/{queued_run_id}/queue-entry")
        assert removed.status_code == 200, removed.text
        assert removed.json() == {"status": "deleted", "id": queued_run_id}
        assert client.get(f"/api/workflow-runs/{queued_run_id}").status_code == 404
        assert client.get(f"/api/documents/{source_document_id}").status_code == 200

        queued_again = client.post(f"/api/workflow-runs/{source_run['id']}/enqueue")
        assert queued_again.status_code == 200, queued_again.text
        queued_again_id = queued_again.json()["id"]
        canceled = client.post(f"/api/workflow-runs/{queued_again_id}/cancel-waiting")
        assert canceled.status_code == 200, canceled.text
        removed_canceled = client.delete(f"/api/workflow-runs/{queued_again_id}/queue-entry")
        assert removed_canceled.status_code == 200, removed_canceled.text
        assert client.get(f"/api/workflow-runs/{queued_again_id}").status_code == 404
        assert client.get(f"/api/documents/{source_document_id}").status_code == 200

        queued_active = client.post(f"/api/workflow-runs/{source_run['id']}/enqueue")
        assert queued_active.status_code == 200, queued_active.text
        queued_active_id = queued_active.json()["id"]
        started = client.post(f"/api/workflow-runs/{queued_active_id}/start")
        assert started.status_code == 200, started.text
        removed_active = client.delete(f"/api/workflow-runs/{queued_active_id}/queue-entry")
        assert removed_active.status_code == 200, removed_active.text
        assert client.get(f"/api/workflow-runs/{queued_active_id}").status_code == 404
        assert client.get(f"/api/documents/{source_document_id}").status_code == 200

        queued_paused = client.post(f"/api/workflow-runs/{source_run['id']}/enqueue")
        assert queued_paused.status_code == 200, queued_paused.text
        queued_paused_id = queued_paused.json()["id"]
        started_paused = client.post(f"/api/workflow-runs/{queued_paused_id}/start")
        assert started_paused.status_code == 200, started_paused.text
        paused = client.post(f"/api/workflow-runs/{queued_paused_id}/pause")
        assert paused.status_code == 200, paused.text
        removed_paused = client.delete(f"/api/workflow-runs/{queued_paused_id}/queue-entry")
        assert removed_paused.status_code == 200, removed_paused.text
        assert client.get(f"/api/workflow-runs/{queued_paused_id}").status_code == 404
        assert client.get(f"/api/documents/{source_document_id}").status_code == 200


def test_workflow_restart_retries_failed_items_without_reupload(monkeypatch) -> None:
    def fake_failed_execute(db, item, graph):
        return {
            "document_id": item.document_id,
            "filename": item.filename,
            "status": "failed",
            "error_message": "VLM_PROVIDER_REQUEST_FAILED: invalid api key",
            "branch_path": None,
            "path_node_ids": ["classifier"],
            "completed_node_ids": [],
            "current_node_id": None,
            "current_node_kind": None,
            "current_node_label": None,
            "node_results": {},
        }

    def fake_success_execute(db, item, graph):
        return {
            "document_id": item.document_id,
            "filename": item.filename,
            "status": "completed",
            "error_message": None,
            "branch_path": "class:contract",
            "path_node_ids": ["classifier"],
            "completed_node_ids": ["classifier"],
            "current_node_id": None,
            "current_node_kind": None,
            "current_node_label": None,
            "node_results": {},
            "classification": {"status": "classified", "class_name": "contract"},
            "kie_values": {},
            "required_overall_status": None,
            "required_items": {},
        }

    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        monkeypatch.setattr("app.workflows._execute_graph_for_item", fake_failed_execute)
        with get_client() as client:
            schema = create_schema(client, name="workflow_restart_schema")
            classifier = create_document_classifier(client, name="workflow_restart_classifier")
            checklist = create_required_field_checklist(client, name="workflow_restart_checklist")
            workflow = client.post(
                "/api/workflows",
                json={"name": "restart_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
            )
            assert workflow.status_code == 200, workflow.text
            run_response = client.post(
                f"/api/workflows/{workflow.json()['id']}/runs",
                files=[
                    ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert run_response.status_code == 200, run_response.text
            run_id = run_response.json()["id"]
            failed_run = client.get(f"/api/workflow-runs/{run_id}").json()
            assert failed_run["status"] == "completed_with_errors"
            assert failed_run["failed_count"] == 2
            document_ids = [item["document_id"] for item in failed_run["items"]]

            monkeypatch.setattr("app.workflows._execute_graph_for_item", fake_success_execute)
            restarted = client.post(f"/api/workflow-runs/{run_id}/restart")
            assert restarted.status_code == 200, restarted.text
            restarted_id = restarted.json()["id"]
            assert restarted_id != run_id
            payload = client.get(f"/api/workflow-runs/{restarted_id}").json()
            assert payload["status"] == "completed"
            assert payload["restarted_from_run_id"] == run_id
            assert payload["uploaded_count"] == 2
            assert [item["document_id"] for item in payload["items"]] == document_ids
            assert {item["status"] for item in payload["items"]} == {"completed"}

            restarted_completed = client.post(f"/api/workflow-runs/{restarted_id}/restart")
            assert restarted_completed.status_code == 200, restarted_completed.text
            second_restart_id = restarted_completed.json()["id"]
            assert second_restart_id != restarted_id
            payload = client.get(f"/api/workflow-runs/{second_restart_id}").json()
            assert payload["status"] == "completed"
            assert [item["document_id"] for item in payload["items"]] == document_ids
            assert {item["status"] for item in payload["items"]} == {"completed"}
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_workflow_branch_without_fallback_exports_classification_only(monkeypatch) -> None:
    monkeypatch.setattr(
        "app.document_modules.classify_document_with_vlm",
        lambda classes, allow_unknown, image_paths: {
            "status": "unknown",
            "class_name": None,
            "confidence": 0.2,
            "reason": "No class matched.",
            "evidence": [],
        },
    )
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="workflow_fallback_schema")
            classifier = create_document_classifier(client, name="workflow_fallback_classifier")
            checklist = create_required_field_checklist(client, name="workflow_fallback_checklist")
            definition = workflow_definition(schema["id"], classifier["id"], checklist["id"])
            definition["edges"] = [edge for edge in definition["edges"] if edge.get("sourceHandle") == "class:contract" or edge["source"] != "branch"]
            workflow = client.post("/api/workflows", json={"name": "fallback_missing", "definition": definition})
            assert workflow.status_code == 200, workflow.text
            assert workflow.json()["validation_warnings"]

            run_response = client.post(
                f"/api/workflows/{workflow.json()['id']}/runs",
                files=[("files", ("invoice.png", ONE_BY_ONE_PNG, "image/png"))],
            )
            assert run_response.status_code == 200, run_response.text
            run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
            assert run["status"] == "completed"
            item = run["items"][0]
            assert item["status"] == "completed"
            assert item["error_message"] is None
            assert item["result"]["branch_path"] == "unknown"
            assert item["result"]["kie_values"] == {}
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_workflow_branch_without_downstream_path_completes_classification_only() -> None:
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="workflow_no_downstream_schema")
            classifier = create_document_classifier(client, name="workflow_no_downstream_classifier")
            checklist = create_required_field_checklist(client, name="workflow_no_downstream_checklist")
            definition = workflow_definition(schema["id"], classifier["id"], checklist["id"])
            definition["edges"] = [edge for edge in definition["edges"] if edge["source"] != "branch"]
            workflow = client.post("/api/workflows", json={"name": "classification_only", "definition": definition})
            assert workflow.status_code == 200, workflow.text
            assert workflow.json()["validation_warnings"]

            run_response = client.post(
                f"/api/workflows/{workflow.json()['id']}/runs",
                files=[("files", ("contract.png", ONE_BY_ONE_PNG, "image/png"))],
            )
            assert run_response.status_code == 200, run_response.text
            run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
            item = run["items"][0]
            assert item["status"] == "completed"
            assert item["error_message"] is None
            assert item["result"]["classification"]["class_name"] == "contract"
            assert item["result"]["branch_path"] == "class:contract"
            assert item["result"]["kie_values"] == {}
            assert item["result"]["required_items"] == {}
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_workflow_progress_payload_marks_current_and_completed_nodes(monkeypatch) -> None:
    observed_current_nodes: list[str | None] = []

    def fake_classifier(db, document_id, node):
        item = db.query(WorkflowRunItem).filter(WorkflowRunItem.document_id == document_id).one()
        observed_current_nodes.append(json.loads(item.result_json).get("current_node_kind"))
        return {
            "kind": "classifier",
            "status": "completed",
            "classification": {"status": "classified", "class_name": "contract"},
        }

    monkeypatch.setattr("app.workflows._execute_classifier_node", fake_classifier)
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="workflow_progress_schema")
            classifier = create_document_classifier(client, name="workflow_progress_classifier")
            checklist = create_required_field_checklist(client, name="workflow_progress_checklist")
            definition = workflow_definition(schema["id"], classifier["id"], checklist["id"])
            definition["edges"] = [edge for edge in definition["edges"] if edge["source"] != "branch"]
            workflow = client.post("/api/workflows", json={"name": "progress_payload", "definition": definition})
            assert workflow.status_code == 200, workflow.text

            run_response = client.post(
                f"/api/workflows/{workflow.json()['id']}/runs",
                files=[("files", ("progress.png", ONE_BY_ONE_PNG, "image/png"))],
            )
            assert run_response.status_code == 200, run_response.text
            run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
            result = run["items"][0]["result"]
            assert observed_current_nodes == ["classifier"]
            assert result["current_node_id"] is None
            assert "classifier" in result["completed_node_ids"]
            assert "branch" in result["completed_node_ids"]
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_workflow_run_uses_workflow_worker_limit_for_blocking_work(monkeypatch) -> None:
    active = 0
    max_active = 0
    lock = threading.Lock()

    def fake_execute(db, item, graph):
        nonlocal active, max_active
        with lock:
            active += 1
            max_active = max(max_active, active)
        time.sleep(0.05)
        with lock:
            active -= 1
        return {
            "document_id": item.document_id,
            "filename": item.filename,
            "status": "completed",
            "error_message": None,
            "branch_path": "class:contract",
            "path_node_ids": ["classifier"],
            "completed_node_ids": ["classifier"],
            "current_node_id": None,
            "current_node_kind": None,
            "current_node_label": None,
            "node_results": {},
            "classification": {"status": "classified", "class_name": "contract"},
            "kie_values": {},
            "required_overall_status": None,
            "required_items": {},
        }

    monkeypatch.setattr("app.workflows._execute_graph_for_item", fake_execute)
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        os.environ["WORKFLOW_MAX_WORKERS"] = "2"
        os.environ["VLM_MAX_CONCURRENT_REQUESTS"] = "8"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="workflow_parallel_schema")
            classifier = create_document_classifier(client, name="workflow_parallel_classifier")
            checklist = create_required_field_checklist(client, name="workflow_parallel_checklist")
            workflow = client.post(
                "/api/workflows",
                json={"name": "parallel_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
            )
            assert workflow.status_code == 200, workflow.text
            run_response = client.post(
                f"/api/workflows/{workflow.json()['id']}/runs",
                files=[("files", (f"doc_{index}.png", ONE_BY_ONE_PNG, "image/png")) for index in range(4)],
            )
            assert run_response.status_code == 200, run_response.text
            run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
            assert run["status"] == "completed"
            assert max_active == 2
    finally:
        os.environ.pop("WORKFLOW_MAX_WORKERS", None)
        os.environ.pop("VLM_MAX_CONCURRENT_REQUESTS", None)
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_workflow_parallel_item_failure_does_not_stop_other_items(monkeypatch) -> None:
    def fake_execute(db, item, graph):
        if item.filename == "doc_1.png":
            raise RuntimeError("synthetic workflow failure")
        return {
            "document_id": item.document_id,
            "filename": item.filename,
            "status": "completed",
            "error_message": None,
            "branch_path": "class:contract",
            "path_node_ids": ["classifier"],
            "completed_node_ids": ["classifier"],
            "current_node_id": None,
            "current_node_kind": None,
            "current_node_label": None,
            "node_results": {},
            "classification": {"status": "classified", "class_name": "contract"},
            "kie_values": {},
            "required_overall_status": None,
            "required_items": {},
        }

    monkeypatch.setattr("app.workflows._execute_graph_for_item", fake_execute)
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        os.environ["VLM_MAX_CONCURRENT_REQUESTS"] = "3"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="workflow_failure_schema")
            classifier = create_document_classifier(client, name="workflow_failure_classifier")
            checklist = create_required_field_checklist(client, name="workflow_failure_checklist")
            workflow = client.post(
                "/api/workflows",
                json={"name": "parallel_failure", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
            )
            assert workflow.status_code == 200, workflow.text
            run_response = client.post(
                f"/api/workflows/{workflow.json()['id']}/runs",
                files=[("files", (f"doc_{index}.png", ONE_BY_ONE_PNG, "image/png")) for index in range(3)],
            )
            assert run_response.status_code == 200, run_response.text
            run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
            statuses = {item["filename"]: item["status"] for item in run["items"]}
            assert run["status"] == "completed_with_errors"
            assert statuses["doc_1.png"] == "failed"
            assert statuses["doc_0.png"] == "completed"
            assert statuses["doc_2.png"] == "completed"
    finally:
        os.environ.pop("VLM_MAX_CONCURRENT_REQUESTS", None)
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_workflow_retry_failed_items_requeues_only_failures(monkeypatch) -> None:
    attempts: dict[str, int] = {}

    def fake_execute(db, item, graph):
        attempts[item.filename] = attempts.get(item.filename, 0) + 1
        if item.filename == "doc_1.png" and attempts[item.filename] == 1:
            raise RuntimeError("synthetic transient failure")
        return {
            "document_id": item.document_id,
            "filename": item.filename,
            "status": "completed",
            "error_message": None,
            "branch_path": "class:contract",
            "path_node_ids": ["classifier"],
            "completed_node_ids": ["classifier"],
            "current_node_id": None,
            "current_node_kind": None,
            "current_node_label": None,
            "node_results": {},
            "classification": {"status": "classified", "class_name": "contract"},
            "kie_values": {},
            "required_overall_status": None,
            "required_items": {},
        }

    monkeypatch.setattr("app.workflows._execute_graph_for_item", fake_execute)
    try:
        os.environ["VLM_PROVIDER"] = "mock"
        os.environ["VLM_MAX_CONCURRENT_REQUESTS"] = "2"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client, name="workflow_retry_schema")
            classifier = create_document_classifier(client, name="workflow_retry_classifier")
            checklist = create_required_field_checklist(client, name="workflow_retry_checklist")
            workflow = client.post(
                "/api/workflows",
                json={"name": "retry_failed_workflow", "definition": workflow_definition(schema["id"], classifier["id"], checklist["id"])},
            )
            assert workflow.status_code == 200, workflow.text
            run_response = client.post(
                f"/api/workflows/{workflow.json()['id']}/runs",
                files=[("files", (f"doc_{index}.png", ONE_BY_ONE_PNG, "image/png")) for index in range(2)],
            )
            assert run_response.status_code == 200, run_response.text
            run = client.get(f"/api/workflow-runs/{run_response.json()['id']}").json()
            assert run["status"] == "completed_with_errors"
            assert run["failed_count"] == 1

            retry = client.post(f"/api/workflow-runs/{run['id']}/retry-failed")
            assert retry.status_code == 200, retry.text
            retried = client.get(f"/api/workflow-runs/{run['id']}").json()
            assert retried["status"] == "completed"
            assert retried["completed_count"] == 2
            assert retried["failed_count"] == 0
            assert attempts == {"doc_0.png": 1, "doc_1.png": 2}
    finally:
        os.environ.pop("VLM_MAX_CONCURRENT_REQUESTS", None)
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_batch_worker_exception_does_not_leave_batch_running(monkeypatch) -> None:
    from app import extraction as extraction_module

    monkeypatch.setattr("app.main.run_batch_jobs", lambda batch_id, job_ids: None)

    try:
        os.environ["VLM_PROVIDER"] = "mock"
        get_settings.cache_clear()
        with get_client() as client:
            schema = create_schema(client)
            response = client.post(
                "/api/batches",
                data={"schema_id": schema["id"]},
                files=[
                    ("files", ("first.png", ONE_BY_ONE_PNG, "image/png")),
                    ("files", ("second.png", ONE_BY_ONE_PNG, "image/png")),
                ],
            )
            assert response.status_code == 200, response.text
            batch = response.json()
            job_ids = [item["job_id"] for item in batch["items"]]
            original_run_job = extraction_module.run_extraction_job

            def flaky_run_job(job_id: str) -> None:
                if job_id == job_ids[0]:
                    raise RuntimeError("worker boom")
                original_run_job(job_id)

            monkeypatch.setattr(extraction_module, "run_extraction_job", flaky_run_job)
            extraction_module.run_batch_jobs(batch["id"], job_ids)

            loaded = client.get(f"/api/batches/{batch['id']}")
            assert loaded.status_code == 200, loaded.text
            payload = loaded.json()
            assert payload["status"] == "completed_with_errors"
            assert payload["progress"] == 1
            assert payload["failed_count"] == 1
            assert payload["completed_count"] == 1
            assert {item["status"] for item in payload["items"]} == {"completed", "failed"}
    finally:
        os.environ["VLM_PROVIDER"] = "openai"
        get_settings.cache_clear()


def test_raw_dependency_imports() -> None:
    import bleach  # noqa: F401
    import mammoth  # noqa: F401
    import openpyxl  # noqa: F401
    import pptx  # noqa: F401


def test_raw_extraction_pdf_upload() -> None:
    with get_client() as client:
        response = client.post(
            "/api/raw-extractions",
            files={"file": ("sample.pdf", make_pdf_bytes(), "application/pdf")},
        )
        assert response.status_code == 200, response.text
        payload = response.json()
        assert payload["status"] == "completed"
        assert payload["source_format"] == "pdf"
        assert payload["pdf_url"]
        assert payload["html_url"]

        pdf = client.get(payload["pdf_url"])
        assert pdf.status_code == 200
        assert pdf.headers["content-type"] == "application/pdf"

        html_response = client.get(payload["html_url"])
        assert html_response.status_code == 200
        assert "Invoice No. INV-2026-001" in html_response.text

        recent = client.get("/api/raw-extractions").json()
        assert any(item["id"] == payload["id"] for item in recent)


def test_raw_extraction_pdf_upload_with_images_option() -> None:
    with get_client() as client:
        response = client.post(
            "/api/raw-extractions",
            data={"include_images": "true"},
            files={"file": ("sample.pdf", make_pdf_with_image_bytes(), "application/pdf")},
        )
        assert response.status_code == 200, response.text
        payload = response.json()
        assert payload["status"] == "completed"

        html_response = client.get(payload["html_url"])
        assert html_response.status_code == 200
        assert "data:image/png;base64" in html_response.text


def test_raw_extraction_pdf_upload_includes_images_by_default() -> None:
    with get_client() as client:
        response = client.post(
            "/api/raw-extractions",
            files={"file": ("sample.pdf", make_pdf_with_image_bytes(), "application/pdf")},
        )
        assert response.status_code == 200, response.text
        payload = response.json()
        assert payload["status"] == "completed"

        html_response = client.get(payload["html_url"])
        assert html_response.status_code == 200
        assert "data:image/png;base64" in html_response.text


def test_raw_extraction_pptx_upload_with_images_option(monkeypatch) -> None:
    def fake_convert(source_path, suffix, pdf_path):
        document = fitz.open()
        page = document.new_page(width=240, height=120)
        page.insert_text((24, 60), f"Preview for {source_path.name}")
        document.save(pdf_path)
        document.close()

    monkeypatch.setattr("app.raw_extractor.convert_office_to_pdf", fake_convert)

    with get_client() as client:
        response = client.post(
            "/api/raw-extractions",
            data={"include_images": "true"},
            files={
                "file": (
                    "deck_with_image.pptx",
                    make_pptx_with_image_bytes(),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )
        assert response.status_code == 200, response.text
        payload = response.json()
        assert payload["status"] == "completed", payload
        html_response = client.get(payload["html_url"])
        assert html_response.status_code == 200
        assert "data:image/png;base64" in html_response.text


def test_raw_extraction_office_uploads(monkeypatch) -> None:
    def fake_convert(source_path, suffix, pdf_path):
        document = fitz.open()
        page = document.new_page(width=240, height=120)
        page.insert_text((24, 60), f"Preview for {source_path.name}")
        document.save(pdf_path)
        document.close()

    monkeypatch.setattr("app.raw_extractor.convert_office_to_pdf", fake_convert)

    samples = [
        ("report.docx", make_docx_bytes(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "Quarterly Report"),
        ("book.xlsx", make_xlsx_bytes(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "Revenue"),
        ("deck.pptx", make_pptx_bytes(), "application/vnd.openxmlformats-officedocument.presentationml.presentation", "Roadmap"),
    ]
    with get_client() as client:
        for filename, data, mime_type, expected_text in samples:
            response = client.post(
                "/api/raw-extractions",
                files={"file": (filename, data, mime_type)},
            )
            assert response.status_code == 200, response.text
            payload = response.json()
            assert payload["status"] == "completed", payload
            assert payload["pdf_url"]
            assert payload["html_url"]
            assert client.get(payload["pdf_url"]).status_code == 200
            html_response = client.get(payload["html_url"])
            assert html_response.status_code == 200
            assert expected_text in html_response.text


def test_raw_extraction_xlsx_formula_option(monkeypatch) -> None:
    def fake_convert(source_path, suffix, pdf_path):
        document = fitz.open()
        page = document.new_page(width=240, height=120)
        page.insert_text((24, 60), f"Preview for {source_path.name}")
        document.save(pdf_path)
        document.close()

    monkeypatch.setattr("app.raw_extractor.convert_office_to_pdf", fake_convert)

    with get_client() as client:
        response = client.post(
            "/api/raw-extractions",
            data={"include_formulas": "true"},
            files={
                "file": (
                    "book.xlsx",
                    make_xlsx_formula_bytes(),
                    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                )
            },
        )
        assert response.status_code == 200, response.text
        payload = response.json()
        assert payload["status"] == "completed", payload
        html_response = client.get(payload["html_url"])
        assert html_response.status_code == 200
        assert "=SUM(B2:B2)" in html_response.text


def upload_png(client):
    response = client.post(
        "/api/documents",
        files={"file": ("invoice.png", ONE_BY_ONE_PNG, "image/png")},
    )
    assert response.status_code == 200, response.text
    return response.json()


def create_schema(client, name: str | None = None):
    global SCHEMA_COUNTER
    if name is None:
        SCHEMA_COUNTER += 1
        schema_name = "invoice_basic" if SCHEMA_COUNTER == 1 else f"invoice_basic_{SCHEMA_COUNTER}"
    else:
        schema_name = name
    response = client.post(
        "/api/schemas",
        json={
            "name": schema_name,
            "display_name": schema_name.replace("_", " ").title(),
            "fields": [
                {
                    "key_name": "invoice_number",
                    "description": "Invoice number near the top of the document. Return null if missing.",
                    "output_format": "string",
                },
                {
                    "key_name": "total_amount",
                    "description": "Final total amount including tax.",
                    "output_format": "float",
                },
            ],
        },
    )
    assert response.status_code == 200, response.text
    return response.json()


def create_document_classifier(client, name: str = "document_classifier"):
    response = client.post(
        "/api/document-classifiers",
        json={
            "name": name,
            "description": "문서를 사용자가 정의한 class 후보 중 하나로 분류합니다.",
            "allow_unknown": True,
            "classes": [
                {
                    "class_name": "contract",
                    "description": "계약 조건과 서명 또는 날인이 있는 문서",
                    "signals": ["계약", "서명", "날인"],
                },
                {
                    "class_name": "consent_form",
                    "description": "개인정보 또는 금융정보 조회 동의 여부가 있는 문서",
                    "signals": ["동의", "개인정보", "조회"],
                },
            ],
        },
    )
    assert response.status_code == 200, response.text
    return response.json()


def create_required_field_checklist(client, name: str = "required_checklist"):
    response = client.post(
        "/api/required-field-checklists",
        json={
            "name": name,
            "description": "필수 항목의 존재 여부만 확인합니다.",
            "regions": [
                {"id": "signature_region", "name": "서명 영역", "page": 1, "x": 0.55, "y": 0.55, "width": 0.35, "height": 0.25}
            ],
            "items": [
                {
                    "item_name": "성명",
                    "description": "성명이 문서에 존재하는지 확인합니다.",
                    "evidence_type": "text_or_handwriting",
                    "required": True,
                },
                {
                    "item_name": "서명",
                    "description": "서명 또는 날인이 존재하는지 확인합니다.",
                    "evidence_type": "signature_or_stamp",
                    "required": True,
                    "region_id": "signature_region",
                },
                {
                    "item_name": "체크박스",
                    "description": "필수 체크박스 표시가 존재하는지 확인합니다.",
                    "evidence_type": "checkbox",
                    "required": True,
                },
            ],
        },
    )
    assert response.status_code == 200, response.text
    return response.json()


def workflow_definition(schema_id: str, classifier_id: str, checklist_id: str):
    return {
        "nodes": [
            {"id": "input", "position": {"x": 0, "y": 120}, "data": {"kind": "input", "label": "Input"}},
            {
                "id": "classifier",
                "position": {"x": 220, "y": 120},
                "data": {"kind": "classifier", "label": "Classifier", "config": {"classifier_id": classifier_id}},
            },
            {"id": "branch", "position": {"x": 440, "y": 120}, "data": {"kind": "branch", "label": "Branch"}},
            {
                "id": "kie_contract",
                "position": {"x": 680, "y": 60},
                "data": {"kind": "kie", "label": "Contract KIE", "config": {"schema_id": schema_id}},
            },
            {
                "id": "required_contract",
                "position": {"x": 900, "y": 60},
                "data": {"kind": "required-checker", "label": "Contract Required", "config": {"checklist_id": checklist_id}},
            },
            {"id": "merge", "position": {"x": 1120, "y": 120}, "data": {"kind": "merge", "label": "Merge"}},
            {"id": "export", "position": {"x": 1340, "y": 120}, "data": {"kind": "export", "label": "Export"}},
        ],
        "edges": [
            {"id": "input-classifier", "source": "input", "target": "classifier"},
            {"id": "classifier-branch", "source": "classifier", "target": "branch"},
            {"id": "branch-contract", "source": "branch", "target": "kie_contract", "sourceHandle": "class:contract"},
            {"id": "branch-consent-form", "source": "branch", "target": "merge", "sourceHandle": "class:consent_form"},
            {"id": "branch-unknown", "source": "branch", "target": "merge", "sourceHandle": "unknown"},
            {"id": "kie-required", "source": "kie_contract", "target": "required_contract"},
            {"id": "required-merge", "source": "required_contract", "target": "merge"},
            {"id": "merge-export", "source": "merge", "target": "export"},
        ],
    }


def make_pdf_bytes() -> bytes:
    document = fitz.open()
    page = document.new_page(width=240, height=120)
    page.insert_text((24, 60), "Invoice No. INV-2026-001")
    buffer = io.BytesIO()
    document.save(buffer)
    document.close()
    return buffer.getvalue()


def make_pdf_with_image_bytes() -> bytes:
    document = fitz.open()
    page = document.new_page(width=240, height=120)
    page.insert_text((24, 28), "Document with image")
    page.insert_image(fitz.Rect(24, 40, 80, 96), stream=ONE_BY_ONE_PNG)
    buffer = io.BytesIO()
    document.save(buffer)
    document.close()
    return buffer.getvalue()


def make_docx_bytes() -> bytes:
    buffer = io.BytesIO()
    document_xml = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    <w:p><w:pPr><w:pStyle w:val="Heading1"/></w:pPr><w:r><w:t>Quarterly Report</w:t></w:r></w:p>
    <w:p><w:r><w:t>Executive summary paragraph.</w:t></w:r></w:p>
    <w:tbl>
      <w:tr><w:tc><w:p><w:r><w:t>Metric</w:t></w:r></w:p></w:tc><w:tc><w:p><w:r><w:t>Value</w:t></w:r></w:p></w:tc></w:tr>
      <w:tr><w:tc><w:p><w:r><w:t>Revenue</w:t></w:r></w:p></w:tc><w:tc><w:p><w:r><w:t>100</w:t></w:r></w:p></w:tc></w:tr>
    </w:tbl>
  </w:body>
</w:document>"""
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>""",
        )
        archive.writestr("word/document.xml", document_xml)
    return buffer.getvalue()


def make_xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Finance"
    sheet.append(["Metric", "Value"])
    sheet.append(["Revenue", 100])
    buffer = io.BytesIO()
    workbook.save(buffer)
    workbook.close()
    return buffer.getvalue()


def make_xlsx_formula_bytes() -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Finance"
    sheet.append(["Metric", "Value"])
    sheet.append(["Revenue", 100])
    sheet.append(["Total", "=SUM(B2:B2)"])
    buffer = io.BytesIO()
    workbook.save(buffer)
    workbook.close()
    return buffer.getvalue()


def make_pptx_bytes() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Roadmap"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(6), Inches(1))
    textbox.text_frame.text = "Launch Raw Data Extractor"
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def make_pptx_with_image_bytes() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Image slide"
    image_stream = io.BytesIO(ONE_BY_ONE_PNG)
    slide.shapes.add_picture(image_stream, Inches(1), Inches(1.4), width=Inches(1))
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()
